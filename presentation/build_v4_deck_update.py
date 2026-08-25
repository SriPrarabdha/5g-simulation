"""Add the v4 optimizer evidence to CDOT_UPF_Steering_Evidence_Review.pptx.

The deck is hand-composed on Blank layouts, so this script reproduces its
design tokens rather than using placeholders: sampled colours, the same
Georgia/Segoe UI/Consolas roles, the same 0.72in margin, footer rule, section
tag and page-number furniture.

Two slides go into the opening section (headline plus caveats, because the
caveats have to be read before any number is), and six into the optimizer
section. A handful of existing lines that the v4 campaign has since falsified
are corrected in place; everything else in the deck is left untouched.
"""
from __future__ import annotations

import copy
import json
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Emu, Inches, Pt

ROOT = Path(__file__).resolve().parent
DECK = ROOT / "CDOT_UPF_Steering_Evidence_Review.pptx"
FIG = ROOT / "generated_assets" / "v4"
DATA = json.load(open(FIG / "deck-data.json"))

# ---- design tokens sampled from the existing deck -------------------------
BG        = RGBColor(0xFA, 0xFA, 0xFB)
CARD      = RGBColor(0xFF, 0xFF, 0xFF)
BORDER    = RGBColor(0xDE, 0xE3, 0xE8)
RULE      = RGBColor(0xEB, 0xEE, 0xF1)
TEAL      = RGBColor(0x0F, 0x4C, 0x5C)
INK       = RGBColor(0x11, 0x17, 0x21)
BODY      = RGBColor(0x3C, 0x46, 0x53)
MUTED     = RGBColor(0x6C, 0x76, 0x83)
FAINT     = RGBColor(0x98, 0xA2, 0xAE)
GREEN     = RGBColor(0x2E, 0x7D, 0x5B)
RED       = RGBColor(0x9C, 0x3B, 0x35)
AMBER     = RGBColor(0xB5, 0x76, 0x1F)
SOFT      = RGBColor(0xEC, 0xF3, 0xF0)

SERIF, SANS, MONO = "Georgia", "Segoe UI", "Consolas"
M = 0.72               # page margin
W = 11.89              # content width


def _tb(slide, l, t, w, h, text, *, font=SANS, size=12.0, bold=False,
        color=BODY, align=PP_ALIGN.LEFT, spacing=1.0, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = spacing
        r = p.add_run()
        r.text = line
        r.font.name = font
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
    return box


def _rect(slide, l, t, w, h, fill, line=None, line_w=0.6):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(line_w)
    sh.shadow.inherit = False
    return sh


def chrome(slide, eyebrow, title, section, *, intro=None, intro_w=9.28):
    """Background, eyebrow, title, optional intro, footer rule and section tag."""
    _rect(slide, 0, 0, 13.33, 7.5, BG)
    _tb(slide, M, 0.52, W, 0.22, eyebrow, font=MONO, size=10.5, bold=True, color=TEAL)
    _tb(slide, M, 0.80, W, 0.90, title, font=SERIF, size=30, color=INK)
    if intro:
        _tb(slide, M, 1.60, intro_w, 0.60, intro, size=13, color=BODY, spacing=1.25)
    _rect(slide, M, 6.98, W, 0.01, RULE)
    _tb(slide, M, 7.08, 7.14, 0.22, section, font=MONO, size=8.5, color=FAINT)


def card(slide, l, t, w, h, label, body, *, accent=False, label_color=MUTED,
         body_size=12.4, body_color=BODY):
    _rect(slide, l, t, w, h, CARD, BORDER)
    if accent:
        _rect(slide, l, t, 0.04, h, TEAL)
    _tb(slide, l + 0.22, t + 0.18, w - 0.44, 0.22, label, font=MONO, size=9.19,
        bold=True, color=label_color)
    _tb(slide, l + 0.22, t + 0.52, w - 0.44, h - 0.70, body, size=body_size,
        color=body_color, spacing=1.22)


def stat(slide, l, t, w, value, caption, *, color=TEAL, size=30):
    """The big-number tile used on the deck's summary slides."""
    _rect(slide, l, t, w, 1.28, CARD, BORDER)
    _tb(slide, l + 0.18, t + 0.16, w - 0.36, 0.50, value, font=SERIF, size=size, color=color)
    _tb(slide, l + 0.18, t + 0.78, w - 0.36, 0.42, caption, font=MONO, size=8.2,
        bold=True, color=MUTED, spacing=1.18)


def bullet(slide, l, t, w, head, body, *, gap=0.27, body_h=0.62):
    _rect(slide, l, t + 0.07, 0.07, 0.07, TEAL)
    _tb(slide, l + 0.20, t, w - 0.20, 0.28, head, size=11.4, bold=True, color=INK)
    _tb(slide, l + 0.20, t + gap, w - 0.20, body_h, body, size=10.7, color=BODY, spacing=1.20)


def picture(slide, name, l, t, w):
    return slide.shapes.add_picture(str(FIG / name), Inches(l), Inches(t), width=Inches(w))


# ---------------------------------------------------------------------------
def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def move_slide(prs, frm, to):
    xml_slides = prs.slides._sldIdLst
    ids = list(xml_slides)
    xml_slides.remove(ids[frm])
    xml_slides.insert(to, ids[frm])


def renumber(prs):
    """Page numbers are baked into a right-aligned mono textbox on each slide."""
    for i, s in enumerate(prs.slides):
        for sh in s.shapes:
            if not sh.has_text_frame or sh.left is None:
                continue
            if sh.left > Inches(11.5) and sh.top > Inches(6.9):
                t = sh.text_frame.text.strip()
                if t.isdigit():
                    p = sh.text_frame.paragraphs[0]
                    if p.runs:
                        p.runs[0].text = f"{i + 1:02d}"
                        for extra in p.runs[1:]:
                            extra.text = ""
    return prs


def pagenum(slide):
    _tb(slide, 11.73, 7.04, 0.88, 0.22, "00", font=MONO, size=8.5,
        color=FAINT, align=PP_ALIGN.RIGHT)


def replace_text(slide, old, new):
    """Swap a run's text while keeping its formatting."""
    for sh in slide.shapes:
        if not sh.has_text_frame:
            continue
        for p in sh.text_frame.paragraphs:
            joined = "".join(r.text for r in p.runs)
            if old in joined and p.runs:
                p.runs[0].text = joined.replace(old, new)
                for extra in p.runs[1:]:
                    extra.text = ""
                return True
    return False


# ===================== new slides: opening section =========================
def slide_headline(prs):
    s = blank(prs)
    chrome(s, "THE V4 RESULT", "Predictive steering wins where the loss is declared",
           "00 · THE BIG PICTURE",
           intro="Re-run after an audit of the evaluator, on 25 000 paired simulated days. "
                 "The controllers now beat the simple rule — but only for capacity loss that is "
                 "announced in advance, and the boundary is as important as the gain.",
           intro_w=11.4)
    tiles = [("+24.0%", "OVERLOAD REMOVED\nON UNSEEN SEEDS", TEAL),
             ("0 / 473", "REGRESSIONS ON\nDECLARED EVENTS", GREEN),
             ("36 / 160", "CONFIGURATIONS PASSED\nEVERY GATE", GREEN),
             ("3 / 4", "SURVIVED FRESH-SEED\nVALIDATION", GREEN),
             ("25 000", "PAIRED SIMULATED\nDAYS", MUTED)]
    w, gap = 2.21, 0.20
    for i, (v, c, col) in enumerate(tiles):
        stat(s, M + i * (w + gap), 2.42, w, v, c, color=col, size=26)
    card(s, M, 3.98, 5.78, 2.62, "WHAT THE CONTROLLERS NOW WIN",
         "On declared maintenance the guarded pre-drain removes 55.2% of uplink overload, and on "
         "maintenance followed by a stadium-scale surge, 52.2%. Across 473 pairs of those two "
         "families it did not lose a single one.\n\n"
         "Cohort MPC clears every gate too, at 33.1% and 32.4%, with a tenth of the collateral "
         "harm and a tenth of the routing churn.")
    card(s, 6.83, 3.98, 5.78, 2.62, "WHAT IT STILL COSTS",
         "When maintenance is announced and then a second, unannounced brownout lands on a box the "
         "controller has already drained onto, pre-drain loses in 30.5% of cases. The family is "
         "still net positive, but the worst cases are severe.\n\n"
         "MPC cuts that exposure to 6.6% of cases. That is the trade C-DOT chooses between.")
    pagenum(s)
    return s


def slide_caveats(prs):
    s = blank(prs)
    chrome(s, "READ THIS FIRST", "Five caveats that bound every number in this deck",
           "00 · THE BIG PICTURE",
           intro="None of these are disclaimers added afterwards. Each one changes how a figure on "
                 "a later slide should be read.", intro_w=11.4)
    items = [
        ("This is a shadow controller in simulation.",
         "No policy touched live traffic. Every record in the project is labelled synthetic and none "
         "of it is calibrated to real C-DOT capacity figures."),
        ("Only about 3% of load is controllable.",
         "The lever is where NEW sessions are placed. Sessions already established are never "
         "migrated, so a controller can only steer the margin."),
        ("Gains are for announced capacity loss only.",
         "With no declared event the controller publishes bit-exact static output. Both pure-surprise "
         "families are deliberate ties, verified byte for byte."),
        ("One stress family produced no measurable stress.",
         "An unannounced demand surge never overloaded any box, even at 4x arrivals across a zone. "
         "That family scores a guaranteed zero, so the headline is conservative."),
        ("Chosen on one seed pool, judged on another.",
         "The headline numbers come from a pool never inspected before the candidates were frozen. "
         "Screening numbers are 3-4 points higher and are not the claim."),
    ]
    for i, (head, body) in enumerate(items):
        col, row = divmod(i, 3)
        l = M if col == 0 else 6.83
        t = 2.40 + row * 1.44
        _rect(s, l, t, 5.78, 1.30, CARD, BORDER)
        _rect(s, l, t, 0.04, 1.30, TEAL)
        _tb(s, l + 0.26, t + 0.15, 5.30, 0.20, f"CAVEAT {i+1}", font=MONO, size=8.6,
            bold=True, color=MUTED)
        _tb(s, l + 0.26, t + 0.42, 5.30, 0.26, head, size=11.4, bold=True, color=INK)
        _tb(s, l + 0.26, t + 0.72, 5.30, 0.50, body, size=10.2, color=BODY, spacing=1.18)
    _rect(s, 6.83, 5.28, 5.78, 1.30, CARD, BORDER)
    _rect(s, 6.83, 5.28, 0.04, 1.30, AMBER)
    _tb(s, 7.09, 5.43, 5.30, 0.20, "HOW TO READ THE OPTIMIZER SECTION", font=MONO, size=8.6,
        bold=True, color=AMBER)
    _tb(s, 7.09, 5.72, 5.30, 0.26, "Experiments 1-6 are the earlier campaign.", size=11.4,
        bold=True, color=INK)
    _tb(s, 7.09, 6.02, 5.30, 0.50, "They remain on the record for method and for the negative "
        "results. Experiment 7 supersedes their conclusion.", size=10.2, color=BODY, spacing=1.18)
    pagenum(s)
    return s


# ===================== new slides: optimizer section ========================
def slide_audit(prs):
    s = blank(prs)
    chrome(s, "EXPERIMENT 7", "Re-running the campaign after auditing the evaluator",
           "03 · THE OPTIMIZER",
           intro="The earlier campaign concluded no controller could beat the simple rule. An audit "
                 "found three defects — all in the scoring code and the scenario generator, none in "
                 "the controllers themselves. Fixing them changed the answer.", intro_w=11.4)
    cards = [
        ("DEFECT 1 · UNITS", "The optimiser and the scoreboard disagreed",
         "The simulator scores overload as load ÷ capacity − 1, a relative quantity. The safety "
         "guard and the pre-drain solver both worked in absolute Mbps. On a box cut to a tenth of "
         "its envelope those differ tenfold, so the controller undervalued exactly the action it "
         "exists to take."),
        ("DEFECT 2 · METRIC", "One stress family scored infinity",
         "An unavailable box has zero capacity, so those pairs scored infinity for both controllers "
         "— an exact tie that still failed a finite-metric check on all 160 configurations. That "
         "single check made the contract unpassable before any controller ran."),
        ("DEFECT 3 · SIZING", "The MPC horizon was set in windows, not hours",
         "Halving the control cadence quintupled the optimisation, so every 2-minute configuration "
         "hit a 2-second solver budget that the acceptance rules allowed to be 120 seconds. The "
         "reported cadence finding was measuring a sizing bug."),
    ]
    w = 3.83
    for i, (label, head, body) in enumerate(cards):
        l = M + i * (w + 0.20)
        card(s, l, 2.62, w, 2.72, label, "", accent=True)
        _tb(s, l + 0.26, 2.94, w - 0.50, 0.52, head, size=12.6, bold=True, color=INK, spacing=1.10)
        _tb(s, l + 0.26, 3.56, w - 0.50, 1.62, body, size=10.6, color=BODY, spacing=1.20)
    _rect(s, M, 5.62, W, 1.00, SOFT)
    _rect(s, M, 5.62, 0.04, 1.00, TEAL)
    _tb(s, M + 0.26, 5.78, W - 0.52, 0.22, "WHAT CHANGED AS A RESULT", font=MONO, size=9.5,
        bold=True, color=TEAL)
    _tb(s, M + 0.26, 6.06, W - 0.52, 0.45,
        "Configurations passing every gate: 0 of 160 → 36 of 160.   Worst stress family: −1.7% → "
        "+24.2%.   MPC solver timeouts: 40% at 10-minute cadence and 100% at 2-minute → zero.",
        size=11.4, color=BODY, spacing=1.16)
    pagenum(s)
    return s


def band(slide, label, body, *, t=5.86, accent=TEAL):
    """Bottom strip used to close a chart slide, matching the deck's callout."""
    _rect(slide, M, t, W, 0.80, SOFT)
    _rect(slide, M, t, 0.04, 0.80, accent)
    _tb(slide, M + 0.26, t + 0.13, W - 0.52, 0.20, label, font=MONO, size=9.5, bold=True, color=accent)
    _tb(slide, M + 0.26, t + 0.38, W - 0.52, 0.42, body, size=10.8, color=BODY, spacing=1.18)


def slide_discovery(prs):
    s = blank(prs)
    d = DATA["discovery"]
    chrome(s, "EXPERIMENT 7 · DISCOVERY", "160 configurations, 20 000 paired simulated days",
           "03 · THE OPTIMIZER",
           intro="Every configuration is one point. To pass it must clear a 10% minimum gain and add "
                 "no more than a quarter of the overload it removes — plus nine further checks.",
           intro_w=8.6)
    picture(s, "fig_frontier.png", M, 2.32, 7.60)
    notes = [("36 of 160 cleared every check",
              "against 0 of 160 in the earlier campaign, on the same simulator and the same "
              "24-hour scenarios."),
             ("Pre-drain reaches further",
              "the strongest passing configuration removes 27.7% of overload while adding 9% of what "
              "it removes."),
             ("MPC sits low and left",
              "less gain, but collateral harm near zero — the conservative end of the same frontier.")]
    t = 2.52
    for head, body in notes:
        bullet(s, 8.58, t, 4.03, head, body, body_h=0.70)
        t += 1.12
        if t < 5.6:
            _rect(s, 8.58, t - 0.20, 4.03, 0.01, RULE)
    band(s, 'THE SHAPE OF THE CAMPAIGN',
         'Each configuration ran the same 125 paired days: five stress families x 25 seeds, one per CPU. 20 000 paired 24-hour simulations in total, 160 of 160 arms valid, zero non-finite results.')
    pagenum(s)
    return s


def slide_validation(prs):
    s = blank(prs)
    chrome(s, "EXPERIMENT 7 · VALIDATION", "Frozen candidates, on seeds never inspected",
           "03 · THE OPTIMIZER",
           intro="Configurations that look good on the data that chose them are not evidence. The four "
                 "leaders were frozen by a rule written down in advance — ranked on the lower "
                 "confidence bound, not the headline — then re-run on a separate seed pool.",
           intro_w=8.6)
    picture(s, "fig_shrink.png", M, 2.42, 7.60)
    notes = [("Three of four held every gate",
              "arm 143 fell just below the 10% bar at 9.8% and was not promoted."),
             ("The drop is 3-4 points, not a collapse",
              "that gap is the selection bias being removed. A candidate that fell apart here would "
              "have been a screening artifact; none did."),
             ("Seeds are provably disjoint",
              "from the screening pool, from every earlier campaign, and from the protected test set.")]
    t = 2.52
    for head, body in notes:
        bullet(s, 8.58, t, 4.03, head, body, body_h=0.70)
        t += 1.12
        if t < 5.6:
            _rect(s, 8.58, t - 0.20, 4.03, 0.01, RULE)
    band(s, 'WHY THE SECOND POOL EXISTS',
         '1 250 fresh paired days per candidate, 250 per stress family, drawn from seeds 81 000 and above — disjoint from the 80 000-80 124 screening pool, from every earlier campaign, and from the protected test set.')
    pagenum(s)
    return s


def slide_families(prs):
    s = blank(prs)
    chrome(s, "WHERE IT WINS AND WHERE IT LOSES", "Three states, not two",
           "03 · THE OPTIMIZER",
           intro="A claim that predictive steering is simply better would not survive scrutiny, and it "
                 "is not what the data says. Bars right of the line are overload removed; bars left "
                 "are overload the controller added. Both panels share one scale.", intro_w=11.4)
    picture(s, "fig_families.png", M, 2.30, 11.30)
    band = [("Declared loss — wins, never loses", "473 informative pairs, 0 regressions. Median case "
             "removes 72% (pre-drain) and 40% (MPC) of uplink overload.", GREEN),
            ("Pure surprise — bit-exact ties", "With no declared event the published policy is byte-for-byte "
             "identical to the simple rule. Verified, not assumed.", MUTED),
            ("Declared loss, then a surprise — can lose", "Pre-drain regresses in 74 of 243 pairs; MPC in "
             "16 of 243. Still net positive, but this is the real exposure.", RED)]
    w = 3.83
    for i, (head, body, col) in enumerate(band):
        l = M + i * (w + 0.20)
        _rect(s, l, 5.76, w, 0.92, CARD, BORDER)
        _rect(s, l, 5.76, 0.04, 0.92, col)
        _tb(s, l + 0.24, 5.88, w - 0.46, 0.26, head, size=11.0, bold=True, color=INK)
        _tb(s, l + 0.24, 6.17, w - 0.46, 0.48, body, size=9.9, color=BODY, spacing=1.16)
    pagenum(s)
    return s


def slide_notice(prs):
    s = blank(prs)
    chrome(s, "ADVANCE NOTICE", "The resource that actually matters",
           "03 · THE OPTIMIZER",
           intro="Gain scales with how much warning the maintenance window carries. Notice is a "
                 "property of the scenario, drawn from its seed, so every configuration faced the "
                 "same schedule.", intro_w=8.6)
    picture(s, "fig_notice.png", M, 2.44, 7.30)
    notes = [("Four hours is worth 4x thirty minutes",
              "pre-drain removes 42.2% of overload with four hours' warning and 10.9% with thirty "
              "minutes."),
             ("MPC saturates around 24%",
              "its two-hour planning horizon cannot exploit three or four hours of notice; pre-drain "
              "keeps climbing."),
             ("This is a scheduling lever, not a model",
              "the single cheapest way to increase the benefit is to announce maintenance earlier.")]
    t = 2.52
    for head, body in notes:
        bullet(s, 8.30, t, 4.31, head, body, body_h=0.70)
        t += 1.12
        if t < 5.6:
            _rect(s, 8.30, t - 0.20, 4.31, 0.01, RULE)
    band(s, 'WHAT THIS MEANS OPERATIONALLY',
         'Announcing a maintenance window three hours ahead instead of thirty minutes roughly triples the overload the controller can remove. That is a change to scheduling practice, not to any model.')
    pagenum(s)
    return s


def slide_cadence(prs):
    s = blank(prs)
    c = DATA["cadence"]
    chrome(s, "CADENCE", "Two minutes buys churn, not benefit",
           "03 · THE OPTIMIZER",
           intro="Running the controller every 2 minutes instead of every 10 was tested across all 160 "
                 "configurations, with the optimisation held to the same size so cadence trades "
                 "lookahead against reactivity rather than against solver feasibility.", intro_w=11.4)
    picture(s, "fig_cadence.png", M, 2.28, 10.10)
    rows = [("Pre-drain", "10 min", c["predrain_10"]), ("Pre-drain", "2 min", c["predrain_2"]),
            ("Cohort MPC", "10 min", c["mpc_10"]), ("Cohort MPC", "2 min", c["mpc_2"])]
    lx, ty = 10.98, 2.42
    _tb(s, lx, ty, 1.63, 0.20, "PASSED", font=MONO, size=8.2, bold=True, color=MUTED)
    for i, (name, cad, v) in enumerate(rows):
        y = ty + 0.30 + i * 0.52
        _tb(s, lx, y, 1.63, 0.20, f"{name} · {cad}", font=MONO, size=8.6, color=MUTED)
        col = GREEN if v["passing"] else RED
        _tb(s, lx, y + 0.19, 1.63, 0.26, f"{v['passing']} / {v['n']}", font=SERIF, size=15, color=col)
    _rect(s, M, 5.74, W, 0.88, SOFT)
    _rect(s, M, 5.74, 0.04, 0.88, AMBER)
    _tb(s, M + 0.26, 5.88, W - 0.52, 0.22, "THE SAME CONFIGURATION, RUN AT BOTH CADENCES",
        font=MONO, size=9.5, bold=True, color=AMBER)
    _tb(s, M + 0.26, 6.16, W - 0.52, 0.40,
        "Arm 60 at 10 minutes removes 27.67% of overload with 0.192 churn and passes every check. "
        "Arm 124 — identical apart from cadence — removes 27.69% with 0.702 churn and fails on churn "
        "alone. Faster control bought two hundredths of a point for 3.7x the routing change.",
        size=11.2, color=BODY, spacing=1.16)
    pagenum(s)
    return s


# ===================== corrections to existing slides ======================
# Only lines the v4 campaign has since falsified. Everything else is left alone.
CORRECTIONS = [
    (5, "The platform works. The smart controllers do not.",
        "The platform works. The controllers win where the loss is declared."),
    (5, "Keeping those two findings apart is the point of the whole project.",
        "Superseded by Experiment 7. The finding that matters is which conditions each result holds in."),
    (5, "Pre-drain gained 79–83% on planned outages, lost 5.6–7.2% on surprises.",
        "Pre-drain now removes 55% on declared maintenance and ties exactly on pure surprises."),
    (5, "Break-even in development; 13.3% WORSE on 128 realistic test days.",
        "Now clears every gate at 16.3% on fresh seeds, with near-zero collateral harm."),
    (39, "Deciding where traffic goes — and why it barely helps",
         "Deciding where traffic goes — and when it genuinely helps"),
    (54, "What we run in production, and why",
         "What we recommend running, and why"),
    (54, "The simple rule. Split new sessions across each group's healthy boxes in proportion to "
         "their capacity. No forecast, no memory, no solver.",
         "The simple rule stays the default. Guarded pre-drain or cohort MPC is enabled for declared "
         "maintenance windows only, and falls back to the simple rule everywhere else."),
    (58, "No deployable controller captured it.",
         "Experiment 7 captured part of it, for declared events only."),
    (58, "Under realistic mixed conditions, worst-day limits and speed limits, MPC was break-even to "
         "harmful and pre-drain traded average gain for unsafe worst days. The simple rul",
         "After correcting three evaluator defects, 36 of 160 configurations clear every gate and "
         "three of four survive fresh-seed validation — but only where capacity loss is announced. "
         "Pure surprises remain bit-exact ties."),
]


def main() -> int:
    prs = Presentation(DECK)
    before = len(prs.slides)

    applied, missed = 0, []
    for idx, old, new in CORRECTIONS:
        if replace_text(prs.slides[idx], old, new):
            applied += 1
        else:
            missed.append((idx, old[:52]))

    # opening section: headline then caveats, straight after THE HEADLINE RESULT
    slide_headline(prs); move_slide(prs, len(prs.slides) - 1, 6)
    slide_caveats(prs);  move_slide(prs, len(prs.slides) - 1, 7)

    # optimizer section: after SELF-AUDIT (was 53, now 55), before DECISION
    at = 56
    for builder in (slide_audit, slide_discovery, slide_validation,
                    slide_families, slide_notice, slide_cadence):
        builder(prs)
        move_slide(prs, len(prs.slides) - 1, at)
        at += 1

    renumber(prs)
    prs.save(DECK)

    print(f"slides {before} -> {len(prs.slides)}")
    print(f"corrections applied {applied}/{len(CORRECTIONS)}")
    for idx, frag in missed:
        print(f"  MISSED slide {idx}: {frag!r}")
    return 1 if missed else 0


if __name__ == "__main__":
    raise SystemExit(main())
