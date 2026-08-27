"""Add plain-English bridge slides, presenter notes, and a C-DOT Q&A appendix."""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


DECK = Path(__file__).resolve().parent / "CDOT_UPF_Steering_Evidence_Review.pptx"

BG = RGBColor(0xFA, 0xFA, 0xFB)
CARD = RGBColor(0xFF, 0xFF, 0xFF)
BORDER = RGBColor(0xDE, 0xE3, 0xE8)
RULE = RGBColor(0xEB, 0xEE, 0xF1)
TEAL = RGBColor(0x0F, 0x4C, 0x5C)
INK = RGBColor(0x11, 0x17, 0x21)
BODY = RGBColor(0x3C, 0x46, 0x53)
MUTED = RGBColor(0x6C, 0x76, 0x83)
FAINT = RGBColor(0x98, 0xA2, 0xAE)
GREEN = RGBColor(0x2E, 0x7D, 0x5B)
AMBER = RGBColor(0xB5, 0x76, 0x1F)
SOFT = RGBColor(0xEC, 0xF3, 0xF0)
SERIF, SANS, MONO = "Georgia", "Segoe UI", "Consolas"
M, W = 0.72, 11.89


def tb(slide, left, top, width, height, text, *, font=SANS, size=12, bold=False,
       color=BODY, align=PP_ALIGN.LEFT, spacing=1.0, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    frame = box.text_frame
    frame.word_wrap = True
    frame.margin_left = frame.margin_right = frame.margin_top = frame.margin_bottom = 0
    frame.vertical_anchor = anchor
    for index, line in enumerate(text.split("\n")):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.alignment = align
        paragraph.line_spacing = spacing
        run = paragraph.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return box


def rect(slide, left, top, width, height, fill, line=None):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(0.7)
    shape.shadow.inherit = False
    return shape


def chrome(slide, eyebrow: str, title: str, section: str, intro: str = ""):
    rect(slide, 0, 0, 13.33, 7.5, BG)
    tb(slide, M, 0.52, W, 0.22, eyebrow, font=MONO, size=10.5, bold=True, color=TEAL)
    tb(slide, M, 0.80, W, 0.75, title, font=SERIF, size=29, color=INK)
    if intro:
        tb(slide, M, 1.55, W, 0.52, intro, size=12.5, color=BODY, spacing=1.2)
    rect(slide, M, 6.98, W, 0.01, RULE)
    tb(slide, M, 7.08, 8.0, 0.22, section, font=MONO, size=8.5, color=FAINT)
    tb(slide, 11.73, 7.04, 0.88, 0.22, "00", font=MONO, size=8.5,
       color=FAINT, align=PP_ALIGN.RIGHT)


def card(slide, left, top, width, height, label, heading, body, *, color=TEAL):
    rect(slide, left, top, width, height, CARD, BORDER)
    rect(slide, left, top, 0.05, height, color)
    tb(slide, left + 0.22, top + 0.16, width - 0.44, 0.20, label,
       font=MONO, size=8.8, bold=True, color=color)
    tb(slide, left + 0.22, top + 0.43, width - 0.44, 0.34, heading,
       size=12.0, bold=True, color=INK)
    tb(slide, left + 0.22, top + 0.83, width - 0.44, height - 0.98, body,
       size=10.8, color=BODY, spacing=1.16)


def move_last_slide(presentation: Presentation, to_index: int) -> None:
    items = presentation.slides._sldIdLst
    slide_id = list(items)[-1]
    items.remove(slide_id)
    items.insert(to_index, slide_id)


def titles(presentation: Presentation) -> list[str]:
    result = []
    for slide in presentation.slides:
        text = "\n".join(shape.text for shape in slide.shapes if shape.has_text_frame)
        result.append(text)
    return result


def traffic_bridge(presentation: Presentation) -> None:
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    chrome(
        slide,
        "PLAIN ENGLISH · TRAFFIC MODEL",
        "The mathematics is four familiar ideas stacked together",
        "01 · THE SIMULATOR",
        "The formula on the next slide is bookkeeping, not magic. Read it from left to right.",
    )
    cards = [
        ("01 · BASELINE", "How busy is this service normally?",
         "Each area and service class starts with an average session-arrival rate. This fixes the scale, not the exact future."),
        ("02 · SHAPE", "What time and day is it?",
         "Daily and weekly curves multiply the baseline. Residential traffic peaks in the evening; enterprise traffic peaks during working hours."),
        ("03 · RANDOMNESS", "What makes today different?",
         "Seeded random variation adds persistence, bursts and events. Separate random streams keep arrivals, failures and movement independent."),
        ("04 · SESSIONS", "How long does each decision remain?",
         "Arrivals become sessions with bandwidth and duration. Long-lived sessions create committed load that future decisions cannot undo."),
    ]
    for index, item in enumerate(cards):
        left = M + (index % 2) * 6.05
        top = 2.24 + (index // 2) * 2.10
        card(slide, left, top, 5.82, 1.82, *item)
    rect(slide, M, 6.47, W, 0.32, SOFT)
    tb(slide, M + 0.18, 6.53, W - 0.36, 0.18,
       "OUTPUT CHECK · offered = carried + queued + dropped + rejected. If this identity fails, stop.",
       font=MONO, size=9.3, bold=True, color=GREEN)


def optimizer_bridge(presentation: Presentation) -> None:
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    chrome(
        slide,
        "PLAIN ENGLISH · OPTIMIZER",
        "The optimizer chooses percentages—not individual users",
        "03 · THE OPTIMIZER",
        "For every traffic group it answers one bounded question: what share of new sessions should each eligible UPF receive?",
    )
    card(slide, M, 2.20, 3.72, 3.48, "DECISION VARIABLES", "The numbers it may change",
         "Routing weights such as 20%, 30% and 50%. They must be non-negative and add to 100% for every traffic group.")
    card(slide, 4.82, 2.20, 3.72, 3.48, "OBJECTIVE", "What 'better' means",
         "First avoid overload. Then reduce cross-zone delay and routing churn. Penalties express priorities; they do not remove safety constraints.", color=AMBER)
    card(slide, 8.92, 2.20, 3.69, 3.48, "CONSTRAINTS", "Rules it cannot break",
         "UPF capacity, session limits, eligibility, health, delay and anchored sessions. An independent validator checks the proposed policy again.", color=GREEN)
    rect(slide, M, 5.94, W, 0.72, SOFT)
    tb(slide, M + 0.22, 6.05, W - 0.44, 0.42,
       "The LP looks one step ahead. MPC repeats the same idea across several future steps, acts only on the first, then solves again with fresh data.",
       size=11.4, bold=True, color=TEAL, align=PP_ALIGN.CENTER)


QAS = [
    ("Scope and architecture", [
        ("What exactly can the controller change?", "Only the destination weights for new sessions. It cannot move established sessions in the current contract."),
        ("Why not steer every packet?", "UPF selection is a session-level SMF decision. Packet routing is a different control plane and is outside this project."),
        ("Is session migration impossible?", "Not proven. It is an open C-DOT architecture question. The simulator currently assumes no migration and publishes that limit."),
    ]),
    ("Traffic-model foundations", [
        ("Why simulate sessions instead of packets?", "The steering lever and congestion memory are session-level. Packet simulation at 16 million users would add cost without answering this decision question."),
        ("Why use Poisson arrivals?", "It is a transparent baseline for counts in a short interval. Correlated drift, bursts and events are added separately; operator data should replace the assumptions."),
        ("Why use heavy-tailed session lengths?", "Real workloads contain many short sessions and a few very long ones. Those long sessions are precisely why a bad placement remains committed."),
    ]),
    ("Traffic realism and data", [
        ("Are the UPF capacities real C-DOT values?", "No. They are explicit synthetic envelopes. Real throughput and session limits are required before calibration or a shadow pilot."),
        ("How do you know the simulator is internally correct?", "Conservation identities, deterministic seeds, distribution re-fitting, population conservation, eligibility checks and memory-scaling tests."),
        ("Why is rural scaled larger than urban?", "It is a known synthetic ramp artifact, not a claim about Indian traffic. Real per-area subscriber counts are a priority calibration input."),
    ]),
    ("Forecasting", [
        ("Why keep ridge if LightGBM is more accurate?", "The release rule required every gate. LightGBM improved average accuracy but failed a worst-slice gate; ridge was auditable and its uncertainty coverage was reliable."),
        ("What is a conformal bound?", "An empirical safety margin calibrated on past forecast errors. A p90 bound is designed to cover roughly nine out of ten comparable outcomes."),
        ("How do you prevent the model from seeing the future?", "Every feature has an availability time. Only closed telemetry buckets available at decision time may enter training or inference."),
    ]),
    ("Optimizer mathematics", [
        ("What are the optimizer's unknowns?", "One non-negative routing weight for each eligible traffic-group/UPF pair, plus explicit slack or activation variables in some formulations."),
        ("What does the objective minimize?", "Overload first, then smaller penalties for utilization, delay and policy churn. Safety is also enforced as constraints and checked independently."),
        ("What is the difference between LP and MPC?", "The LP optimizes the present step. MPC plans over a horizon, applies only the first decision and re-plans as telemetry changes."),
    ]),
    ("Controller evidence", [
        ("Why does the oracle remove all overload?", "It knows the future event schedule and solves the whole day jointly. It is an upper bound, not a deployable controller."),
        ("Why does pre-drain help planned maintenance?", "Advance notice lets it stop adding long-lived sessions to a UPF before capacity falls."),
        ("Why can pre-drain lose under mixed stress?", "It may drain sessions onto a UPF that later suffers an unannounced event. Those sessions are then anchored."),
    ]),
    ("Experiment 7 and safety", [
        ("Why did Experiment 7 reverse the earlier conclusion?", "An audit found unit, finite-metric and horizon-sizing defects in the evaluator/scenario machinery. Correcting them made the declared-event benefit measurable."),
        ("Does 0/473 regressions prove universal safety?", "No. It covers the informative declared-event pairs in the stated seed pools. Mixed declared-plus-surprise cases still contain regressions."),
        ("What does +24.0% mean?", "Held-out reduction in the pre-declared uplink-overload metric versus matched static runs. It is not a throughput increase or a live-network result."),
    ]),
    ("Operations and cluster", [
        ("Why use 10-minute control instead of two minutes?", "Two-minute control produced almost no extra gain and about 3.7 times the routing churn for a matched configuration."),
        ("How does the cluster work scale?", "Across independent matched scenario/seed pairs. One simulation is normally one shard; it is not divided across 160 nodes."),
        ("What happens if the solver or telemetry fails?", "The independent gate rejects the recommendation and retains the last safe static policy."),
    ]),
    ("Deployment challenge questions", [
        ("Why not use reinforcement learning?", "RL inherits the same limited steering lever and adds validation difficulty. Settle the lever, telemetry and safety contract first."),
        ("What would you need from C-DOT next?", "The SMF steering key, maintenance-notice workflow, real UPF safe envelopes and telemetry/reset semantics."),
        ("Are you asking us to deploy this controller?", "No. The recommendation is a bounded shadow-advisory replay using real telemetry, with no live policy publication."),
    ]),
]


def qa_cover(presentation: Presentation) -> None:
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    chrome(slide, "APPENDIX · C-DOT DISCUSSION", "Questions you are likely to be asked",
           "APPENDIX · PRESENTER SAFETY NET",
           "Twenty-seven direct answers. Use these slides only when a question comes up; do not present them linearly.")
    rect(slide, M, 2.35, W, 3.55, TEAL)
    tb(slide, 1.10, 2.82, 11.10, 0.56, "A strong answer has three parts", font=SERIF,
       size=25, color=CARD, align=PP_ALIGN.CENTER)
    for index, (number, heading, body) in enumerate([
        ("1", "Answer the exact question", "Lead with yes, no, or the measured result."),
        ("2", "State the evidence boundary", "Synthetic, matched scope, declared event, or open assumption."),
        ("3", "Name the next check", "What real C-DOT input or shadow test resolves the uncertainty."),
    ]):
        left = 1.08 + index * 4.05
        tb(slide, left, 3.64, 0.42, 0.50, number, font=SERIF, size=26, color=RGBColor(0x58, 0xE0, 0xD3))
        tb(slide, left + 0.55, 3.60, 3.12, 0.34, heading, size=12.0, bold=True, color=CARD)
        tb(slide, left + 0.55, 4.03, 3.12, 0.72, body, size=10.5, color=RGBColor(0xD8, 0xEA, 0xEC))
    tb(slide, M, 6.23, W, 0.32, "If you do not know a C-DOT fact, say so. Convert it into a shadow-pilot input—not a guess.",
       size=11.2, bold=True, color=AMBER, align=PP_ALIGN.CENTER)


def qa_slide(presentation: Presentation, category: str, questions: list[tuple[str, str]], index: int) -> None:
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    chrome(slide, f"Q&A · {index:02d}", category, "APPENDIX · C-DOT QUESTIONS")
    for row, (question, answer) in enumerate(questions):
        top = 1.80 + row * 1.68
        rect(slide, M, top, W, 1.46, CARD, BORDER)
        rect(slide, M, top, 0.05, 1.46, TEAL if row != 2 else AMBER)
        tb(slide, M + 0.22, top + 0.16, 0.48, 0.28, f"Q{(index-1)*3+row+1:02d}",
           font=MONO, size=10.0, bold=True, color=TEAL)
        tb(slide, M + 0.76, top + 0.13, 4.20, 0.52, question, size=12.0, bold=True, color=INK)
        tb(slide, 5.12, top + 0.13, 7.20, 1.03, answer, size=11.1, color=BODY, spacing=1.18)
    tb(slide, M, 6.62, W, 0.20, "Answer first · boundary second · next check third", font=MONO,
       size=9.0, bold=True, color=FAINT, align=PP_ALIGN.CENTER)


NOTES = {
    "THE PROBLEM": "PLAIN ENGLISH: We control only where a new session starts. We do not control the load already committed by existing sessions.\n\nSAY: Think of choosing a parking building when a car enters. Once parked, we currently assume the car cannot be moved.\n\nIF ASKED: Session migration is an open C-DOT question, not a proven impossibility.\n\nDO NOT CLAIM: That we steer packets or live traffic.",
    "HOW THE SYSTEM WORKS": "PLAIN ENGLISH: Every ten minutes the system observes closed telemetry, predicts demand, proposes weights, validates them, and either recommends them or retains static.\n\nSAY: The validator—not the optimizer—has the final word.\n\nIF ASKED: A decision at time t affects only sessions beginning after t.",
    "THE V4 RESULT": "PLAIN ENGLISH: The controller helps when capacity loss is announced. It deliberately behaves like static when there is no declared event.\n\nSAY: The boundary is the result: declared loss, yes; pure surprise, tie; declared loss followed by another surprise, residual risk.\n\nDO NOT CLAIM: Universal 24% improvement.",
    "METHOD": "PLAIN ENGLISH: We decide the success rules before looking at protected results and compare controllers on exactly matched days.\n\nSAY: Same seed, same failure, same traffic—only the controller changes.\n\nIF ASKED: Confidence intervals and worst-case days prevent a good average from hiding unsafe cases.",
    "THE ENGINE": "PLAIN ENGLISH: Each 30-second step updates time-of-day traffic, events, session arrivals, placement, capacity and telemetry.\n\nSAY: Separate random-number streams stop a code change in failures from silently changing arrivals.\n\nIF ASKED: This is discrete-time session simulation, not packet simulation.",
    "STRUCTURE": "PLAIN ENGLISH: We divide demand into 96 groups so unlike services are not averaged together.\n\nSAY: A stadium uploader and an IoT meter should not share one forecast or one routing policy.\n\nIF ASKED: The selection key is area/DNN/slice; 5QI is descriptive unless C-DOT confirms otherwise.",
    "AREAS": "PLAIN ENGLISH: Area labels choose service mixes and daily shapes. The current population scale is synthetic.\n\nSAY: Rural appearing largest is a known index-ramp artifact, not a claim about Indian usage.\n\nDO NOT CLAIM: Geographic calibration.",
    "TRAFFIC CLASSES": "PLAIN ENGLISH: Twelve service classes differ in duration, bandwidth direction, latency sensitivity and time-of-day pattern.\n\nSAY: The differences matter because a long IoT or industrial session commits capacity differently from a short social upload.\n\nIF ASKED: 5QI comes from the service description; it is not automatically a UPF-selection key.",
    "PLAIN ENGLISH · TRAFFIC MODEL": "SAY: There are only four layers: average scale, time shape, randomness and session persistence. The next slide writes these four layers as equations.\n\nIF ASKED: Every random draw is seeded and named so a result can be reproduced.",
    "THE MODEL": "PLAIN ENGLISH: Arrival rate equals baseline multiplied by daily/weekly/event factors and correlated noise. Each arrival then receives bandwidth and duration.\n\nSAY: The equations generate offered demand; capacity is applied afterwards.\n\nDO NOT CLAIM: That the parameter values came from C-DOT.",
    "THE MODEL\nThe realism layer": "PLAIN ENGLISH: The first 16-week corpus used the simpler generator. The newer realism layer adds correlated bandwidth and heavy-tailed durations and validates them by re-fitting output.\n\nSAY: The realism layer is ready, but the 112-day history has not yet been regenerated with it.\n\nDO NOT MIX: Results from the older corpus with claims about the newer generator.",
    "STRESS": "PLAIN ENGLISH: Demand surges and capacity failures are independent so a controller cannot learn that one always predicts the other.\n\nSAY: Planned maintenance has notice; a surprise surge intentionally does not.\n\nIF ASKED: Mixed stress is the hard case because pre-drained sessions may be committed to the next failing UPF.",
    "VERIFICATION": "PLAIN ENGLISH: We ask whether the generated output matches the distributions and conservation laws requested in the configuration.\n\nSAY: Internal validity is strong; external calibration to C-DOT is not yet done.\n\nDO NOT CLAIM: Realism merely because the graphs look plausible.",
    "DATA QUALITY": "PLAIN ENGLISH: The simulator creates clean truth and separately creates damaged telemetry. The pipeline must detect gaps, stale values and counter resets.\n\nSAY: Never smooth across a reset; that invents traffic.\n\nIF ASKED: Two training rounds were discarded when the quality aggregation semantics were found wrong.",
    "SCALE": "PLAIN ENGLISH: Long runs stream results rather than accumulating them in memory. Independent scenario/seed shards use cluster parallelism.\n\nSAY: We scale across experiments, not by splitting one ordinary simulation over 160 nodes.",
    "KEY NUMBERS · THROUGHPUT": "PLAIN ENGLISH: These are measured outputs from one frozen synthetic calibration run, not configuration targets.\n\nSAY: Offered is user demand; carried is served traffic; drops and rejects explain the difference.\n\nDO NOT CLAIM: That the totals match C-DOT production traffic.",
    "PERSPECTIVE": "PLAIN ENGLISH: Some dimensions resemble national scale, while geography, capacity and subscriber calibration remain synthetic.\n\nSAY: The honest answer is 'partly'—large enough to stress the machinery, not yet operator-calibrated.",
    "SCALING UP · STAGE 1": "PLAIN ENGLISH: This test checks whether a longer simulation leaks memory.\n\nSAY: Seven times the duration used only 4.1% more peak memory, supporting streaming execution.",
    "SCALING UP · STAGE 2": "PLAIN ENGLISH: This test finds how many independent workers a node can sustain before memory or scheduling overhead harms throughput.\n\nSAY: Packing is measured rather than assumed.",
    "SCALING UP · STAGE 3": "PLAIN ENGLISH: Scale-out is gated: prove the same inputs on two nodes, then four, then twelve.\n\nSAY: Each node owns independent work; matching hashes and complete shard reports gate the next rung.",
    "METHOD\nHow the forecaster works": "PLAIN ENGLISH: The model predicts new-session demand per group at several horizons, plus an uncertainty bound.\n\nSAY: Existing carried load is handled separately because it is already committed.\n\nIF ASKED: Only features available at forecast issue time are permitted.",
    "RESULTS\nAccuracy on 16 weeks of history": "PLAIN ENGLISH: WAPE is total absolute error divided by total actual demand. Lower is better.\n\nSAY: It is volume-weighted and understandable, but it can hide a bad small slice, so worst-slice gates are separate.\n\nDO NOT CLAIM: 7.63% error for every group or every event.",
    "COMPARISON\nAgainst two simple baselines": "PLAIN ENGLISH: A model earns value only by beating strong simple forecasts on identical rows.\n\nSAY: Yesterday and seasonal baselines also obey the no-future rule.\n\nIF ASKED: The comparison is causal and matched.",
    "MODEL COMPARISON": "PLAIN ENGLISH: The most accurate model is not automatically released. Every accuracy, coverage, peak and worst-slice gate must pass.\n\nSAY: LightGBM won average accuracy but lost the release decision.\n\nDO NOT CLAIM: That ridge is the globally best forecasting algorithm.",
    "FAILURE MODE": "PLAIN ENGLISH: Overall averages look acceptable while the event periods that matter remain hard.\n\nSAY: This is why control evidence is required; forecast improvement does not guarantee operational improvement.",
    "COST\nHow long it takes to train and to run": "PLAIN ENGLISH: Training happens offline; inference happens on every decision round and must fit the operational deadline.\n\nSAY: A more complex model is acceptable only if its accuracy benefit survives the end-to-end latency and release gates.",
    "DECISION\nWhich model we use, and why": "PLAIN ENGLISH: Ridge stays because the selection contract says no challenger passed every gate.\n\nSAY: Keeping a simpler model is the consequence of the rule, not personal preference.\n\nIF ASKED: Conformal bounds give the optimizer a calibrated upper-demand estimate.",
    "THE CONSTRAINT": "PLAIN ENGLISH: Only the small flow of new sessions is movable; hours of existing load are anchored.\n\nSAY: This limited lever explains why prediction alone often has modest effect.\n\nIF ASKED: The controllable fraction varies by scenario; the deck summarizes it as about 3%.",
    "PLAIN ENGLISH · OPTIMIZER": "SAY: The variables are percentages, the objective ranks undesirable outcomes, and constraints are hard rules. MPC repeats the same calculation over a horizon.\n\nIF ASKED: A separate policy validator recomputes feasibility before publication.",
    "THE CANDIDATES": "PLAIN ENGLISH: These controllers differ in how much future information and state they use. Static is the reference; oracle is an impossible upper bound.\n\nSAY: Reactive is intuitive but harmful because it reacts after sessions are already committed.\n\nDO NOT CLAIM: Any candidate is deployed on C-DOT.",
    "EXPERIMENT 1": "PLAIN ENGLISH: The LP found a valid mathematical solution but barely improved overload because it could move only new sessions.\n\nSAY: Solver success and operational success are different questions.",
    "EXPERIMENT 2": "PLAIN ENGLISH: Changing penalty weights did not repair the weak lever.\n\nSAY: Tuning a formulation cannot create controllability that the architecture does not provide.",
    "EXPERIMENT 3": "PLAIN ENGLISH: The oracle asks whether improvement is theoretically available if the future were known.\n\nSAY: It proves headroom exists, especially around known failures; it is not deployable evidence.",
    "EXPERIMENT 4": "PLAIN ENGLISH: MPC predicts several future steps, applies one decision and replans. Earlier test pools gave unstable results.\n\nSAY: Keep this negative history because it explains why Experiment 7 required an evaluator audit.\n\nDO NOT CLAIM: The old −13.3% is the current conclusion.",
    "EXPERIMENT 5": "PLAIN ENGLISH: MPC needs to estimate how much placed load will survive into each future step. Kaplan–Meier learns that survival curve from start/stop logs without peeking at simulator settings.\n\nIF ASKED: It handles censored sessions that are still running at the end of an observation period.",
    "EXPERIMENT 6": "PLAIN ENGLISH: Pre-drain uses advance maintenance notice to reduce new placements before capacity drops.\n\nSAY: The older configuration campaign exposed the gain-versus-worst-case trade. Experiment 7 later corrected evaluator defects and superseded the release conclusion.",
    "COST": "PLAIN ENGLISH: End-to-end decision time matters, not only solver time.\n\nSAY: The latency budget includes feature preparation, solve, validation and policy construction.\n\nIF ASKED: Faster cadence can raise churn without meaningful gain.",
    "THE FULL SCORECARD": "PLAIN ENGLISH: A candidate must pass the whole row; one failed safety or evidence check blocks promotion.\n\nSAY: This slide records the earlier campaign. Experiment 7's 36 passing configurations are the current result.",
    "THE MACHINERY": "PLAIN ENGLISH: Proposed actions move through preflight, solve, independent validation and fallback.\n\nSAY: A feasible solver status is necessary but not sufficient for policy publication.",
    "SELF-AUDIT": "PLAIN ENGLISH: The team found that some 'solved' plans still exceeded a constraint slightly.\n\nSAY: The correct response was a new mandatory validation gate, not rounding the violation away.\n\nIF ASKED: The affected candidates were already rejected, so no unsafe release occurred.",
    "EXPERIMENT 7": "PLAIN ENGLISH: Three evaluator/scenario defects—not controller changes—had made the contract effectively impossible to pass.\n\nSAY: Units, infinite tie metrics and horizon sizing were corrected, then candidates were re-run on frozen pools.\n\nDO NOT CLAIM: That auditing always improves results; here it changed the measured answer.",
    "EXPERIMENT 7 · DISCOVERY": "PLAIN ENGLISH: Each dot is one configuration tested on the same 125 paired days. Thirty-six passed every gate.\n\nSAY: Discovery selects candidates; it is not the headline evidence because selection creates optimism.",
    "EXPERIMENT 7 · VALIDATION": "PLAIN ENGLISH: Frozen leaders were tested on seeds never used for selection. Three of four kept every gate.\n\nSAY: The 3–4 point drop is expected selection bias being removed.",
    "WHERE IT WINS": "PLAIN ENGLISH: There are three operational states: declared loss, pure surprise and declared loss followed by surprise.\n\nSAY: Predictive wins the first, ties the second, and has residual risk in the third.\n\nDO NOT CLAIM: 0 regressions across all stress families.",
    "ADVANCE NOTICE": "PLAIN ENGLISH: Notice time is a resource. More warning lets long-lived sessions drain naturally before capacity falls.\n\nSAY: Three hours instead of thirty minutes can matter more than changing the model.",
    "CADENCE": "PLAIN ENGLISH: Solving more often did not materially improve gain, but it changed routing weights much more.\n\nSAY: Two minutes bought 0.02 percentage points for 3.7 times the churn in the matched example.",
    "DECISION\nWhat we recommend running, and why": "PLAIN ENGLISH: Static is the reference default. Predictive candidates are for declared-maintenance shadow evaluation only.\n\nSAY: No live policy publication is authorized.\n\nIF ASKED: MPC is the conservative candidate; pre-drain has higher gain and higher mixed-event exposure.",
    "THE PATH FORWARD": "PLAIN ENGLISH: Three operator facts matter more than another model iteration: migration capability, telemetry contract and the right operational loss metric.\n\nSAY: The meeting should assign owners to concrete inputs for a shadow replay.",
}


def slide_text(slide) -> str:
    return "\n".join(shape.text for shape in slide.shapes if shape.has_text_frame and shape.text.strip())


def note_for(text: str) -> str | None:
    # Prefer longer/more specific keys so generic titles such as METHOD or COST
    # do not capture the wrong section.
    for key in sorted(NOTES, key=len, reverse=True):
        if key in text:
            return NOTES[key]
    return None


def write_notes(presentation: Presentation) -> int:
    count = 0
    for slide in presentation.slides:
        note = note_for(slide_text(slide))
        if note:
            frame = slide.notes_slide.notes_text_frame
            frame.text = note
            count += 1
    return count


def renumber(presentation: Presentation) -> None:
    for index, slide in enumerate(presentation.slides, 1):
        for shape in slide.shapes:
            if not shape.has_text_frame or shape.left is None:
                continue
            if shape.left > Inches(11.5) and shape.top > Inches(6.9):
                value = shape.text.strip()
                if value.isdigit():
                    shape.text_frame.paragraphs[0].runs[0].text = f"{index:02d}"


def main() -> int:
    presentation = Presentation(DECK)
    existing = titles(presentation)
    if not any("PLAIN ENGLISH · TRAFFIC MODEL" in text for text in existing):
        traffic_bridge(presentation)
        move_last_slide(presentation, 15)  # before the original mathematical model slide
        optimizer_bridge(presentation)
        move_last_slide(presentation, 44)  # after THE CONSTRAINT, before THE CANDIDATES
        qa_cover(presentation)
        for index, (category, questions) in enumerate(QAS, 1):
            qa_slide(presentation, category, questions, index)
    notes_count = write_notes(presentation)
    renumber(presentation)
    presentation.save(DECK)
    print(f"slides={len(presentation.slides)} presenter_notes={notes_count} qas={sum(len(q) for _, q in QAS)}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
