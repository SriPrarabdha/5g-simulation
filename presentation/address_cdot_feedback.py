#!/usr/bin/env python3
"""Address C-DOT office feedback in the final evidence-review deck.

The deck uses blank layouts and hand-composed shapes.  This update preserves the
existing visual language, rewrites the requested slides, and inserts five
plain-English bridge slides.  It is intentionally idempotent: bridge slides are
not duplicated when the script is run again.
"""
from __future__ import annotations

import shutil
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent
DECK = ROOT / "CDOT_UPF_Steering_Evidence_Review.pptx"
BACKUP = ROOT / "CDOT_UPF_Steering_Evidence_Review.before_feedback_20260827.pptx"

BG = RGBColor(0xFA, 0xFA, 0xFB)
CARD = RGBColor(0xFF, 0xFF, 0xFF)
BORDER = RGBColor(0xDE, 0xE3, 0xE8)
RULE = RGBColor(0xEB, 0xEE, 0xF1)
TEAL = RGBColor(0x0F, 0x4C, 0x5C)
INK = RGBColor(0x11, 0x17, 0x21)
BODY = RGBColor(0x3C, 0x46, 0x53)
MUTED = RGBColor(0x6C, 0x76, 0x83)
FAINT = RGBColor(0x98, 0xA2, 0xAE)
GREEN = RGBColor(0x2E, 0x7D, 0x5B)
AMBER = RGBColor(0xB5, 0x76, 0x1F)
RED = RGBColor(0x9C, 0x3B, 0x35)
PURPLE = RGBColor(0x67, 0x50, 0xA4)
BLUE = RGBColor(0x2F, 0x69, 0x8A)
SOFT = RGBColor(0xEC, 0xF3, 0xF0)
SOFT_AMBER = RGBColor(0xFA, 0xF0, 0xDE)
SOFT_RED = RGBColor(0xFA, 0xE8, 0xE6)
SOFT_BLUE = RGBColor(0xE8, 0xF0, 0xF5)

SERIF, SANS, MONO = "Georgia", "Segoe UI", "Consolas"
M, W = 0.72, 11.89


def rect(slide, left, top, width, height, fill=CARD, line=None, radius=False):
    kind = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(
        kind, Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(0.7)
    shape.shadow.inherit = False
    return shape


def tb(
    slide,
    left,
    top,
    width,
    height,
    text,
    *,
    font=SANS,
    size=11.0,
    bold=False,
    color=BODY,
    align=PP_ALIGN.LEFT,
    anchor=MSO_ANCHOR.TOP,
    margin=0,
    spacing=1.0,
):
    box = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    frame = box.text_frame
    frame.word_wrap = True
    frame.margin_left = frame.margin_right = Inches(margin)
    frame.margin_top = frame.margin_bottom = Inches(margin)
    frame.vertical_anchor = anchor
    for index, line in enumerate(text.split("\n")):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.alignment = align
        paragraph.line_spacing = spacing
        run = paragraph.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return box


def line(slide, x1, y1, x2, y2, color=RULE, width=1.0):
    connector = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(x1), Inches(y1), Inches(x2), Inches(y2),
    )
    connector.line.color.rgb = color
    connector.line.width = Pt(width)
    return connector


def arrow(slide, left, top, width=0.42, height=0.34, color=TEAL):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW, Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def clear(slide) -> None:
    for shape in list(slide.shapes):
        shape._element.getparent().remove(shape._element)


def chrome(slide, eyebrow, title, section, intro="", *, source=""):
    rect(slide, 0, 0, 13.33, 7.5, BG)
    tb(slide, M, 0.52, W, 0.22, eyebrow, font=MONO, size=10.5, bold=True, color=TEAL)
    tb(slide, M, 0.80, W, 0.72, title, font=SERIF, size=28.5, color=INK)
    if intro:
        tb(slide, M, 1.51, W, 0.56, intro, size=11.3, color=BODY, spacing=1.15)
    rect(slide, M, 6.98, W, 0.01, RULE)
    tb(slide, M, 7.08, 8.9, 0.20, section, font=MONO, size=8.3, color=FAINT)
    if source:
        tb(slide, 7.6, 7.08, 4.35, 0.20, source, font=MONO, size=7.0, color=FAINT,
           align=PP_ALIGN.RIGHT)
    tb(slide, 11.96, 7.04, 0.65, 0.22, "00", font=MONO, size=8.5, color=FAINT,
       align=PP_ALIGN.RIGHT)


def card(slide, left, top, width, height, label, heading, body, *, accent=TEAL,
         body_size=10.2, heading_size=12.0, fill=CARD):
    rect(slide, left, top, width, height, fill, BORDER)
    rect(slide, left, top, 0.045, height, accent)
    tb(slide, left + 0.22, top + 0.16, width - 0.42, 0.20, label,
       font=MONO, size=8.6, bold=True, color=accent)
    tb(slide, left + 0.22, top + 0.43, width - 0.42, 0.38, heading,
       size=heading_size, bold=True, color=INK)
    tb(slide, left + 0.22, top + 0.84, width - 0.42, height - 0.97, body,
       size=body_size, color=BODY, spacing=1.12)


def stat(slide, left, top, width, value, label, *, color=TEAL):
    rect(slide, left, top, width, 1.03, CARD, BORDER)
    tb(slide, left + 0.15, top + 0.13, width - 0.30, 0.38, value,
       font=SERIF, size=22, color=color, align=PP_ALIGN.CENTER)
    tb(slide, left + 0.10, top + 0.58, width - 0.20, 0.32, label,
       font=MONO, size=7.7, bold=True, color=MUTED, align=PP_ALIGN.CENTER)


def slide_text(slide) -> str:
    return "\n".join(
        shape.text for shape in slide.shapes
        if getattr(shape, "has_text_frame", False) and shape.text.strip()
    )


def find_slide(prs, needle: str):
    matches = [slide for slide in prs.slides if needle in slide_text(slide)]
    if not matches:
        raise ValueError(f"slide not found: {needle}")
    return matches[0]


def index_of(prs, target) -> int:
    for index, slide in enumerate(prs.slides):
        if slide == target:
            return index
    raise ValueError("slide is not in presentation")


def move_last_after(prs, anchor) -> None:
    items = prs.slides._sldIdLst
    last = list(items)[-1]
    items.remove(last)
    items.insert(index_of(prs, anchor) + 1, last)


def add_after(prs, anchor_needle: str, unique_title: str, builder) -> None:
    if any(unique_title in slide_text(slide) for slide in prs.slides):
        return
    anchor = find_slide(prs, anchor_needle)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    builder(slide)
    move_last_after(prs, anchor)


def set_existing_text(shape, value: str, *, size=None, bold=None, color=None) -> None:
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    for index, value_line in enumerate(value.split("\n")):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        run = paragraph.add_run()
        run.text = value_line
        if size is not None:
            run.font.size = Pt(size)
        if bold is not None:
            run.font.bold = bold
        if color is not None:
            run.font.color.rgb = color


def rewrite_headline(prs) -> None:
    slide = find_slide(prs, "The platform works. The controllers win where the loss is declared.")
    clear(slide)
    chrome(
        slide,
        "THE HEADLINE RESULT",
        "The platform works. The controllers win where the loss is declared.",
        "00 · THE BIG PICTURE",
        "The platform evidence is strong. Predictive control is useful only inside a declared capacity event; outside that boundary, Static remains the default.",
        source="Forecast phase 2 · survival phase 3 · oracle study",
    )
    tiles = [
        ("0 bytes", "TRAFFIC ACCOUNTING GAP", GREEN),
        ("4.1%", "RSS GROWTH AT 7× DURATION", GREEN),
        ("11.64%", "LOWER FORECAST ERROR", AMBER),
        ("100%", "CLAIRVOYANT FAULT CEILING", PURPLE),
        ("0 / 473", "DECLARED-EVENT REGRESSIONS", GREEN),
    ]
    tile_w, gap = 2.22, 0.19
    for index, (value, label, color) in enumerate(tiles):
        stat(slide, M + index * (tile_w + gap), 2.20, tile_w, value, label, color=color)

    card(
        slide, M, 3.52, 3.75, 2.95,
        "SESSION-LENGTH ESTIMATION",
        "Learned from ordinary logs",
        "Kaplan–Meier uses only session start/stop records and correctly includes sessions still running when the log ends. In 125 blind trials across five hidden duration families, mean calibration error stayed about 2.9–4.2%; stale data always fell back to Static.",
        accent=GREEN, body_size=9.5,
    )
    card(
        slide, 4.79, 3.52, 3.75, 2.95,
        "WHAT 11.64% MEANS",
        "Error fell; accuracy is not 11.64%",
        "Moving-average WAPE was 14.160%; LightGBM reached 12.512%. Relative error reduction = (14.160−12.512)/14.160 = 11.64%. The promotion rule required at least 15% and no unsafe slice, so LightGBM was useful evidence but was not released.",
        accent=AMBER, body_size=9.5,
    )
    card(
        slide, 8.58, 3.52, 4.03, 2.95,
        "CONTROLLERS IN ONE LINE",
        "Now · before · across time",
        "LP: choose the safest traffic split for the next step.\nPre-drain: use a maintenance notice to stop adding sessions to the affected UPF early.\nMPC: plan across a two-hour horizon, apply only the first action, then solve again.\nOracle: knows the future; it is a ceiling, never a deployment option.",
        accent=TEAL, body_size=9.3,
    )


def rewrite_area_slides(prs) -> None:
    area = find_slide(prs, "The eight area types we model")
    for shape in area.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue
        if shape.text.startswith("The simulator does not model geography"):
            set_existing_text(
                shape,
                "Eight archetypes choose service mix and daily shape. Z(k) is a dimensionless multiplier on a service's base arrival rate: below 1 means less synthetic traffic than the reference area; above 1 means more.",
                size=11.0, color=BODY,
            )
        elif shape.text == "AN HONEST LIMITATION OF THE CURRENT BUILD":
            set_existing_text(shape, "HOW TO READ Z(k) — AND WHY IT IS LINEAR", size=9.2, bold=True, color=AMBER)
        elif shape.text.startswith("The area scale is a straight arithmetic ramp"):
            set_existing_text(
                shape,
                "Z(k) multiplies the base session-arrival rate λs. Example: Z=0.70 generates 30% less traffic than the reference Z=1.00; Z=1.40 generates 40% more. The current z(k)=0.70+0.10k ramp is deliberately linear so load rises in equal, auditable steps while we test scaling. It is not a geographic claim: rural is highest only because it is last in the list. Real subscriber counts should replace this ramp during C-DOT calibration.",
                size=9.7, color=BODY,
            )

    daily = find_slide(prs, "How the shape of a day differs by area")
    for shape in daily.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue
        if shape.text.startswith("Each service class carries its own daily curve"):
            set_existing_text(
                shape,
                "A wrapped Gaussian is an ordinary bell-shaped peak placed on a 24-hour clock: distance wraps across midnight, so a 23:30 peak continues smoothly after 00:00. Add one peak for a single rush or two for morning/evening commute. The service class sets this timing; the area sets which classes and how much of each are present.",
                size=10.6, color=BODY,
            )
        elif shape.text == "WHY THE PEAKS MATTER":
            set_existing_text(shape, "WRAPPED GAUSSIAN, IN PLAIN ENGLISH", size=9.2, bold=True, color=TEAL)
        elif shape.text.startswith("The business curve bottoms out"):
            set_existing_text(
                shape,
                "Gaussian = a smooth bell, controlled by centre time c and width w. Wrapped = use circular clock distance δ(h,c)=min(|h−c|, 24−|h−c|), so midnight has no artificial jump. Sum = add multiple bells when a service has more than one busy period. A baseline α keeps traffic non-zero between peaks.",
                size=9.7, color=BODY,
            )
        elif shape.text == "WHAT MULTIPLIES IT":
            set_existing_text(shape, "WHAT HAPPENS AFTER THE DAILY SHAPE", size=9.2, bold=True, color=TEAL)
        elif shape.text.startswith("On top of the daily curve sit three multipliers"):
            set_existing_text(
                shape,
                "The daily curve is multiplied by a weekend factor, seeded weekly noise U[0.86,1.16], and any active surge. This separation makes the cause of a peak explicit: calendar behaviour, week-to-week variation, or an event.",
                size=9.7, color=BODY,
            )


def equation_explainer(slide) -> None:
    chrome(
        slide,
        "THE MODEL · EQUATION-BY-EQUATION",
        "Every equation has one job in the traffic generator",
        "01 · THE SIMULATOR",
        "Read top to bottom: set the scale, shape the clock, draw sessions, keep them alive, then add their load on each UPF.",
        source="build_extreme_history_manifest.py · traffic-model/1.0",
    )
    rows = [
        ("01", "BASE SCALE", "λg = λs · z(k) · S", "Service rate × area multiplier × national scale. This sets expected volume before time-of-day effects.", TEAL),
        ("02", "TIME + EVENTS", "ag(t) = Ds(h) · Ws(d) · νg,w · M(t)", "Daily curve × weekend effect × weekly noise × active surges. Multiplication keeps each cause visible and independently testable.", BLUE),
        ("03", "ARRIVAL COUNT", "Ng,t ~ Poisson(λg · ag(t))", "Turns the expected rate into an integer number of new sessions in this 30-second tick. Poisson is the transparent count baseline; the seed makes it reproducible.", PURPLE),
        ("04", "PERSISTENCE", "Lg ~ DiscreteUniform[Lmin(s), Lmax(s)]", "Each admitted session gets a lifetime. It remains anchored for that many ticks, which is why an early placement can still matter hours later.", AMBER),
        ("05", "BOX LOAD", "RuUL(t) = Σg ng,u(t) · rsUL", "For each UPF, add the uplink rate of every active session it carries. Downlink and session count are accumulated the same way.", GREEN),
        ("06", "SANITY CHECK", "E[A] = λ̄ · E[L]", "Little's Law predicts average active sessions from arrival rate × mean lifetime. Measured occupancy was within 0.04% of this analytic check.", RED),
    ]
    for index, (num, label, formula, body, accent) in enumerate(rows):
        col, row = index % 2, index // 2
        left = M + col * 6.06
        top = 2.22 + row * 1.43
        rect(slide, left, top, 5.82, 1.20, CARD, BORDER)
        rect(slide, left, top, 0.06, 1.20, accent)
        tb(slide, left + 0.20, top + 0.13, 0.38, 0.25, num, font=MONO, size=9, bold=True, color=accent)
        tb(slide, left + 0.64, top + 0.12, 1.25, 0.22, label, font=MONO, size=8.3, bold=True, color=accent)
        tb(slide, left + 1.92, top + 0.09, 3.62, 0.34, formula, font=SERIF, size=14.2, color=INK)
        tb(slide, left + 0.64, top + 0.47, 4.90, 0.59, body, size=9.2, color=BODY, spacing=1.08)
    rect(slide, M, 6.58, W, 0.25, SOFT)
    tb(slide, M + 0.15, 6.61, W - 0.30, 0.18,
       "OUTPUT IDENTITY · offered = carried + queued + dropped + rejected. If it fails, the run is invalid.",
       font=MONO, size=8.6, bold=True, color=GREEN, align=PP_ALIGN.CENTER)


def comparison_slide(slide) -> None:
    chrome(
        slide,
        "POSITIONING",
        "Different simulators answer different 5G questions",
        "01 · THE SIMULATOR",
        "This project is a long-horizon UPF-steering digital twin. Protocol load generators and radio simulators provide different fidelity; a serious validation plan uses them together.",
        source="Public product/project pages checked 27 Aug 2026",
    )
    headers = [
        ("THIS PROJECT", TEAL),
        ("MOBILEUM dsTest", PURPLE),
        ("PACKETRUSHER / UERANSIM", BLUE),
        ("5G-LENA / Simu5G", GREEN),
    ]
    bodies = [
        ("UPF steering over weeks", "Session-level demand, anchored load, faults, forecasts, LP/MPC/pre-drain, paired safety gates.", "Best for: policy evidence and controllability limits."),
        ("Core/CNF lab assurance", "Emulates subscribers, network interfaces and protocol/application flows; stresses real systems under test and reports KPIs.", "Best for: conformance, interoperability and capacity."),
        ("Open UE/gNB core testing", "Exercises registration, sessions, N1/N2 and user-plane traffic against a 5G core with virtual UEs and gNBs.", "Best for: open functional and load tests."),
        ("RAN / NR system simulation", "Models radio, PHY/MAC, scheduling, interference, mobility and packet/data-plane behaviour in ns-3 or OMNeT++.", "Best for: radio and end-to-end research."),
    ]
    gap, width = 0.18, 2.84
    for index, ((header, color), (heading, body, best)) in enumerate(zip(headers, bodies)):
        left = M + index * (width + gap)
        rect(slide, left, 2.25, width, 3.67, CARD, BORDER)
        rect(slide, left, 2.25, width, 0.11, color)
        tb(slide, left + 0.18, 2.52, width - 0.36, 0.25, header,
           font=MONO, size=8.6, bold=True, color=color, align=PP_ALIGN.CENTER)
        tb(slide, left + 0.18, 2.91, width - 0.36, 0.46, heading,
           size=12.2, bold=True, color=INK, align=PP_ALIGN.CENTER)
        tb(slide, left + 0.20, 3.54, width - 0.40, 1.26, body,
           size=9.5, color=BODY, align=PP_ALIGN.CENTER, spacing=1.12)
        line(slide, left + 0.25, 4.92, left + width - 0.25, 4.92, RULE, 0.8)
        tb(slide, left + 0.20, 5.08, width - 0.40, 0.60, best,
           size=9.2, bold=True, color=color, align=PP_ALIGN.CENTER)
    rect(slide, M, 6.15, W, 0.55, SOFT_AMBER)
    tb(slide, M + 0.20, 6.26, W - 0.40, 0.30,
       "HONEST BOUNDARY · We do not model NR radio or emit standards-compliant N1/N2/N3 traffic. Use dsTest or open UE/gNB tools to replay shortlisted steering policies against a real core.",
       size=9.7, bold=True, color=AMBER, align=PP_ALIGN.CENTER)


def runtime_slide(slide) -> None:
    chrome(
        slide,
        "SCALING UP · WALL TIME",
        "A 16-week extreme shard is a ~9.5-hour job—not an 85-minute job",
        "01 · THE SIMULATOR",
        "The 16-week number is a projection from the measured full one-day writer. The 12-node result is a separate measured campaign of many one-day shards.",
        source="extreme-training-runbook.md · production metrics.json",
    )
    rect(slide, M, 2.22, 7.20, 1.22, CARD, BORDER)
    tb(slide, 0.98, 2.42, 1.42, 0.21, "MEASURED", font=MONO, size=8.5, bold=True, color=GREEN)
    tb(slide, 1.76, 2.31, 1.42, 0.46, "5:06.31", font=SERIF, size=25, color=INK, align=PP_ALIGN.CENTER)
    tb(slide, 1.68, 2.83, 1.58, 0.22, "ONE SIMULATED DAY", font=MONO, size=7.8, bold=True, color=MUTED, align=PP_ALIGN.CENTER)
    arrow(slide, 3.35, 2.57, 0.48, 0.34, TEAL)
    tb(slide, 3.94, 2.31, 1.12, 0.46, "× 112", font=SERIF, size=25, color=TEAL, align=PP_ALIGN.CENTER)
    tb(slide, 3.78, 2.83, 1.44, 0.22, "LINEAR PROJECTION", font=MONO, size=7.5, bold=True, color=MUTED, align=PP_ALIGN.CENTER)
    arrow(slide, 5.30, 2.57, 0.48, 0.34, TEAL)
    tb(slide, 5.90, 2.31, 1.58, 0.46, "9:31:46", font=SERIF, size=25, color=PURPLE, align=PP_ALIGN.CENTER)
    tb(slide, 5.78, 2.83, 1.82, 0.22, "16-WEEK PROJECTION", font=MONO, size=7.8, bold=True, color=MUTED, align=PP_ALIGN.CENTER)

    card(slide, M, 3.72, 3.73, 2.62, "ONE SHARD", "More nodes do not shorten it",
         "The current engine is single-shard and CPU-bound. One 112-day simulation stays on one node; assigning 12 nodes to that same shard does not make it 12× faster. Plan a 12-hour reservation until an end-to-end 16-week capacity pilot is measured.",
         accent=PURPLE, body_size=9.6)
    card(slide, 4.66, 3.72, 3.73, 2.62, "MULTIPLE NODES", "Parallelise independent shards",
         "Use different nodes for different seeds, scenarios or controllers. Conservative planning: one extreme 16-week shard per node. Twelve nodes can therefore run twelve independent 16-week shards in one ~12-hour reservation; more shards need another wave.",
         accent=TEAL, body_size=9.6)
    card(slide, 8.60, 2.22, 4.01, 4.12, "WHAT WAS ACTUALLY MEASURED", "384 one-day shards · 12 nodes",
         "Wall time: 5,122.5 s = 85.4 min\nAggregate CPU efficiency: 90.9%\nWorker failures: 0\nSwap: 0\n\nThis proves campaign scale-out and deterministic orchestration. It does not measure a 16-week shard, so 85.4 minutes must not be quoted as the 16-week runtime.",
         accent=GREEN, body_size=9.7)


def model_primer(slide) -> None:
    chrome(
        slide,
        "FORECASTING · FIRST PRINCIPLES",
        "Five model families: what each one is trying to learn",
        "02 · THE FORECASTER",
        "A regression model predicts a number, not a class. Every family sees only information available at issue time; the difference is how flexibly it maps those inputs to future demand.",
        source="forecasting/candidates.py · phase-2 selection",
    )
    rect(slide, M, 2.17, W, 0.62, SOFT_BLUE, BORDER)
    tb(slide, M + 0.18, 2.28, 2.20, 0.20, "REFERENCE · MOVING AVERAGE", font=MONO, size=8.5, bold=True, color=BLUE)
    tb(slide, 2.92, 2.24, 9.42, 0.28,
       "Average the last six completed windows. It is cheap, causal and surprisingly strong—so a complex model must beat it, not a weak straw man.",
       size=10.0, color=BODY)
    items = [
        ("CALENDAR RIDGE", "Stable linear relation", "Nine lag, trend and calendar features are multiplied by learned coefficients. A ridge penalty shrinks unstable coefficients. Simple, fast, auditable; misses nonlinear thresholds.", TEAL),
        ("RIDGE v2", "Same idea, richer causal history", "Adds longer lags, rolling mean/std/max/slope, residual, quality and event features. Still linear, so interactions must be expressed through inputs.", BLUE),
        ("HISTOGRAM GRADIENT", "Many small decision trees", "Sequential trees correct earlier errors and capture thresholds such as 'surge score high AND evening'. Directly fits p50/p90/p95; more flexible, less transparent.", PURPLE),
        ("LIGHTGBM QUANTILE", "Efficient boosted trees", "A highly optimised tree booster for nonlinear interactions. Best pooled WAPE here, but its worst regime/horizon slice regressed beyond the safety guardrail.", GREEN),
        ("REGIME ENSEMBLE", "Switch model by operating state", "Use the normal model ordinarily; switch after a scheduled event or observable surge. It caught peaks best, but overreacted badly in ordinary slices.", AMBER),
    ]
    for index, (label, heading, body, accent) in enumerate(items):
        if index < 3:
            left, top, width = M + index * 4.03, 3.04, 3.83
        else:
            left, top, width = 2.72 + (index - 3) * 4.10, 4.84, 3.90
        card(slide, left, top, width, 1.56, label, heading, body,
             accent=accent, body_size=8.8, heading_size=10.8)


def rewrite_forecaster_flow(prs) -> None:
    slide = find_slide(prs, "How the forecaster works")
    clear(slide)
    chrome(
        slide,
        "METHOD · FORECASTER I/O",
        "Past telemetry goes in; demand distributions come out",
        "02 · THE FORECASTER",
        "One direct model is trained for every group × target × horizon. A direct model predicts that horizon in one step, so errors do not accumulate through recursive forecasts.",
        source="forecasting/bundle.py · candidates.py",
    )
    # Input column
    tb(slide, 0.82, 2.18, 2.40, 0.22, "INPUTS AVAILABLE AT ISSUE TIME", font=MONO, size=8.4, bold=True, color=TEAL, align=PP_ALIGN.CENTER)
    inputs = [
        ("CLOSED 10-MIN BUCKETS", "new sessions · UL · DL"),
        ("CALENDAR", "time/day · daily/weekly cycle"),
        ("OPERATIONS", "declared events · lead time"),
        ("QUALITY", "age · missing · counter reset"),
    ]
    for index, (head, body) in enumerate(inputs):
        top = 2.55 + index * 0.77
        rect(slide, 0.82, top, 2.40, 0.60, CARD, BORDER, radius=True)
        tb(slide, 0.96, top + 0.10, 2.12, 0.18, head, font=MONO, size=7.8, bold=True, color=TEAL, align=PP_ALIGN.CENTER)
        tb(slide, 0.96, top + 0.31, 2.12, 0.18, body, size=8.7, color=BODY, align=PP_ALIGN.CENTER)

    arrow(slide, 3.36, 3.73, 0.50, 0.38, TEAL)
    rect(slide, 3.96, 2.55, 2.66, 3.02, SOFT_BLUE, BORDER, radius=True)
    tb(slide, 4.18, 2.78, 2.22, 0.24, "CAUSAL FEATURE ROW", font=MONO, size=8.6, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
    tb(slide, 4.18, 3.22, 2.22, 1.84,
       "Current ridge\nlast value · 6-window mean · trend\nyesterday · clock/day cycles\n\nChallengers add\nlong lags · rolling volatility\nresidual · event · quality",
       size=9.2, color=BODY, align=PP_ALIGN.CENTER, spacing=1.12)
    tb(slide, 4.18, 5.20, 2.22, 0.20, "NO FUTURE VALUES ALLOWED", font=MONO, size=7.6, bold=True, color=RED, align=PP_ALIGN.CENTER)

    arrow(slide, 6.74, 3.73, 0.50, 0.38, BLUE)
    rect(slide, 7.34, 2.55, 2.10, 3.02, CARD, BORDER, radius=True)
    tb(slide, 7.56, 2.78, 1.66, 0.24, "DIRECT MODELS", font=MONO, size=8.6, bold=True, color=PURPLE, align=PP_ALIGN.CENTER)
    tb(slide, 7.56, 3.21, 1.66, 1.58,
       "96 groups\n× 3 targets\n× 8 horizons\n\n= 2,304 models",
       font=SERIF, size=14, color=INK, align=PP_ALIGN.CENTER, spacing=1.10)
    tb(slide, 7.56, 5.06, 1.66, 0.28, "10–80 min ahead", font=MONO, size=8.2, bold=True, color=PURPLE, align=PP_ALIGN.CENTER)

    arrow(slide, 9.56, 3.73, 0.50, 0.38, PURPLE)
    tb(slide, 10.17, 2.18, 2.24, 0.22, "OUTPUT FORECAST BUNDLE", font=MONO, size=8.4, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
    outputs = [
        ("p50 / p90 / p95", "new sessions"),
        ("p50 / p90 / p95", "new UL Mbps"),
        ("p50 / p90 / p95", "new DL Mbps"),
        ("PROVENANCE", "group · horizon · issued_at · flags"),
    ]
    for index, (head, body) in enumerate(outputs):
        top = 2.55 + index * 0.77
        rect(slide, 10.17, top, 2.24, 0.60, CARD, BORDER, radius=True)
        tb(slide, 10.31, top + 0.10, 1.96, 0.18, head, font=MONO, size=7.8, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
        tb(slide, 10.31, top + 0.31, 1.96, 0.18, body, size=8.7, color=BODY, align=PP_ALIGN.CENTER)
    rect(slide, M, 6.07, W, 0.60, SOFT)
    tb(slide, M + 0.18, 6.17, W - 0.36, 0.32,
       "GROUP = one area × one service class (96 total).  HORIZON = how far ahead the target window begins.  p50 = central estimate; p90/p95 = upper planning bounds—not confidence scores.",
       size=9.5, bold=True, color=TEAL, align=PP_ALIGN.CENTER)


def scoring_slide(slide) -> None:
    chrome(
        slide,
        "FORECASTING · SCORECARD",
        "How error, safety coverage and release gates are calculated",
        "02 · THE FORECASTER",
        "There is no single 'accuracy score'. Point error, upper-bound reliability, event peaks and worst slices answer different operational questions.",
        source="aggregate_forecast_selection.py · metrics.py",
    )
    cards = [
        (M, 2.18, 3.72, 1.62, "POINT ERROR · LOWER IS BETTER", "WAPE = Σ|actual−p50| / Σactual", "A 12.512% WAPE means total absolute error equals 12.512% of total actual volume. It does not promise every row is within 12.512%.", TEAL),
        (4.82, 2.18, 3.72, 1.62, "GAIN VS BASELINE · HIGHER IS BETTER", "(14.160−12.512) / 14.160 = 11.64%", "This is relative error reduction—not 11.64% accuracy. The frozen promotion gate was ≥15% reduction.", AMBER),
        (8.92, 2.18, 3.69, 1.62, "UPPER-BOUND COVERAGE", "p90 = p50 + 90th-percentile past error", "The added margin is the 90th percentile of |actual−p50| on held-out calibration rows. On unseen rows, coverage = fraction(actual≤p90). Below 88% is unsafe; above 95% is overly wide.", PURPLE),
        (M, 4.07, 5.82, 1.56, "EVENT PEAK UNDERPREDICTION", "miss = max(0, actual peak − predicted peak)", "Compare total miss with the baseline during scheduled or causally detected events. Gate: reduce these misses by ≥20%, because missing the peak is riskier than overshooting it.", RED),
        (6.79, 4.07, 5.82, 1.56, "NO SLICE MAY GET >5% WORSE", "slice = one regime or one forecast horizon", "Compute candidate WAPE ÷ baseline WAPE − 1 for each aggregate slice. The maximum must be ≤5%, so a good average cannot hide an unsafe event state or look-ahead time.", GREEN),
    ]
    for left, top, width, height, label, heading, body, accent in cards:
        card(slide, left, top, width, height, label, heading, body,
             accent=accent, body_size=9.1, heading_size=11.0)
    rect(slide, M, 5.92, W, 0.78, SOFT_AMBER)
    tb(slide, M + 0.20, 6.02, W - 0.40, 0.48,
       "WHY LIGHTGBM WAS NOT PROMOTED · WAPE improved 11.64% (<15% gate) and peak misses improved 27.54% (pass), but its worst regime/horizon slice was 14.70% worse (>5% limit). There is no proof that 15% is impossible—only that this frozen data and feature set did not achieve it safely.",
       size=9.6, bold=True, color=AMBER, align=PP_ALIGN.CENTER)


def rewrite_forecast_comparison(prs) -> None:
    slide = find_slide(prs, "Five model families, head to head")
    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue
        if shape.text.startswith("We set the bar before running anything"):
            set_existing_text(
                shape,
                "Promotion required every gate at once, with thresholds frozen before selection: ≥15% relative WAPE reduction, observed p90 coverage between 88–95%, ≥20% lower event peak miss, and ≤5% regression in every aggregate regime or horizon. Freezing the rules prevents lowering the bar after seeing a favourite model.",
                size=10.3, color=BODY,
            )
        elif shape.text == "The two boosting models won":
            set_existing_text(shape, "Best pooled error, but not release-ready", size=10.8, bold=True, color=INK)
        elif shape.text.startswith("and still missed the 15% bar"):
            set_existing_text(shape, "LightGBM: WAPE 14.160% → 12.512% = 11.64% lower error, not 11.64% accuracy. It missed the ≥15% promotion gate.", size=9.6, color=BODY)
        elif shape.text.startswith("The regime-switching model is the interesting failure"):
            set_existing_text(shape, "Why one average is not enough", size=10.8, bold=True, color=INK)
        elif shape.text.startswith("It detects a surge and switches models"):
            set_existing_text(shape, "The regime ensemble cut event peak misses by 64.2%, yet overall WAPE was 20.2% worse and its worst slice was 601.7% worse. Optimising the rare event damaged ordinary periods.", size=9.6, color=BODY)
        elif shape.text == "Tuning hard for the rare event":
            set_existing_text(shape, "Decision: no challenger promoted", size=10.8, bold=True, color=INK)
        elif shape.text == "wrecked the ordinary case.":
            set_existing_text(shape, "All five models failed at least one precommitted gate. Calendar ridge remains the auditable current model; it is not claimed to be globally best.", size=9.6, color=BODY)

    table = find_slide(prs, "Every family, every check")
    replacements = {
        "ERROR": "WAPE ↓",
        "VS BASELINE": "ERROR GAIN",
        "BOUND HOLDS": "p90 COVERAGE",
        "PEAK UNDERPRED.": "PEAK-MISS GAIN",
        "WORST SLICE": "MAX REGRESSION",
        "RELEASED?": "PROMOTED?",
    }
    for shape in table.shapes:
        if getattr(shape, "has_text_frame", False) and shape.text in replacements:
            set_existing_text(shape, replacements[shape.text], size=7.6, bold=True, color=MUTED)

    decision = find_slide(prs, "Which model we use, and why")
    for shape in decision.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue
        if shape.text.startswith("We run the calendar ridge model with conformal bounds"):
            set_existing_text(
                shape,
                "We retain calendar ridge with conformal upper bounds. In the common five-family comparison, LightGBM WAPE was 12.512% versus ridge at 12.729%—only 0.217 percentage points lower—and LightGBM failed the worst-slice guardrail.",
                size=10.6, color=BODY,
            )

    accuracy = find_slide(prs, "Accuracy on 16 weeks of history")
    for shape in accuracy.shapes:
        if getattr(shape, "has_text_frame", False) and shape.text.startswith("The main measure is WAPE"):
            set_existing_text(
                shape,
                "WAPE is total absolute point error divided by total actual volume; lower is better. This 7.63% result belongs to the frozen calendar-ridge bundle's original holdout. The later five-family selection used a separate controlled evaluation, so 7.63% and 12.512% must not be compared directly.",
                size=10.4, color=BODY,
            )


def rewrite_forecast_next_steps(prs) -> None:
    slide = find_slide(prs, "NEXT STEPS\nForecasting")
    clear(slide)
    chrome(
        slide,
        "NEXT STEPS",
        "Forecasting: fix the evaluation, then improve the model",
        "02 · THE FORECASTER",
        "Horizon = how far ahead we predict (10–80 minutes). Group = one area × one service class (96 total). A model-per-group-and-horizon lets a stadium uploader behave differently from rural IoT at 10 versus 80 minutes.",
        source="control_science_v1.json · forecast backlog",
    )
    left_items = [
        ("01", "Use the intended time split", "Train on weeks 1–11, calibrate on 12–13, test on 14–16. Keep every event wholly inside one split."),
        ("02", "Report every operating slice", "Publish normal, surge, reduced-capacity, offline and delay results by horizon—not only a national average."),
        ("03", "Separate sessions, UL and DL", "Give each session its own bandwidth so UL/DL forecasts are not just rescaled copies of session count."),
    ]
    right_items = [
        ("04", "Add weekly and event memory", "Use a true one-week lag, rolling volatility and only those event features known at forecast issue time."),
        ("05", "Select per group × horizon", "Choose ridge, seasonal or LightGBM on validation data for each task; never choose on test data."),
        ("06", "Pool sparse groups carefully", "Share signal across related services while retaining area and service identity; small groups lack enough events alone."),
    ]
    for col, items in enumerate((left_items, right_items)):
        left = M + col * 6.16
        tb(slide, left, 2.19, 5.74, 0.22,
           "EVALUATION FIRST" if col == 0 else "MODEL IMPROVEMENT SECOND",
           font=MONO, size=9.0, bold=True, color=TEAL if col == 0 else PURPLE)
        for row, (num, head, body) in enumerate(items):
            top = 2.57 + row * 1.17
            rect(slide, left, top, 5.74, 0.98, CARD, BORDER)
            tb(slide, left + 0.18, top + 0.16, 0.38, 0.28, num, font=SERIF, size=18,
               color=TEAL if col == 0 else PURPLE, align=PP_ALIGN.CENTER)
            tb(slide, left + 0.72, top + 0.13, 4.76, 0.25, head, size=10.8, bold=True, color=INK)
            tb(slide, left + 0.72, top + 0.43, 4.76, 0.43, body, size=9.2, color=BODY)
    rect(slide, M, 6.18, W, 0.52, SOFT_AMBER)
    tb(slide, M + 0.18, 6.28, W - 0.36, 0.29,
       "SURGE BOUND IS THE PRIORITY · observed coverage during surges was the binding weakness. Better event labels and faster bound adaptation matter before neural models or GPUs.",
       size=9.5, bold=True, color=AMBER, align=PP_ALIGN.CENTER)


def rewrite_optimizer_flow(prs) -> None:
    slide = find_slide(prs, "The optimizer chooses percentages—not individual users")
    clear(slide)
    chrome(
        slide,
        "PLAIN ENGLISH · OPTIMIZER I/O",
        "The optimizer turns demand and capacity into routing percentages",
        "03 · THE OPTIMIZER",
        "It never chooses a named subscriber. For each traffic group, it chooses what share of future session starts should go to each eligible UPF.",
        source="optimization/highs.py · cohort_mpc.py · steering/gate.py",
    )
    # Main flow
    blocks = [
        (0.78, 2.36, 2.48, "INPUTS", "p50/p90/p95 demand\ncurrent anchored load\nUPF capacity + health\neligibility + delay\nmaintenance notice", TEAL),
        (3.72, 2.36, 2.64, "LP SOLVE", "Variables: routing weights\n\nObjective: minimise overload first, then delay + churn\n\nConstraints: weights ≥0, sum=100%, capacity, health, eligibility", PURPLE),
        (6.82, 2.36, 2.42, "CANDIDATE PLAN", "group → UPF weights\npredicted utilisation\nslack / overload\nsolver status + time\nreason + provenance", BLUE),
        (9.70, 2.36, 2.82, "INDEPENDENT GATE", "Recompute feasibility\ncompare with Static\ncheck telemetry age\ncheck latency + churn\n\nPASS → recommendation\nFAIL → last safe Static", GREEN),
    ]
    for left, top, width, label, body, accent in blocks:
        rect(slide, left, top, width, 2.58, CARD, BORDER, radius=True)
        rect(slide, left, top, width, 0.09, accent)
        tb(slide, left + 0.18, top + 0.24, width - 0.36, 0.24, label,
           font=MONO, size=8.8, bold=True, color=accent, align=PP_ALIGN.CENTER)
        tb(slide, left + 0.20, top + 0.67, width - 0.40, 1.68, body,
           size=9.2, color=BODY, align=PP_ALIGN.CENTER, spacing=1.12)
    for x in (3.30, 6.42, 9.30):
        arrow(slide, x, 3.42, 0.34, 0.30, TEAL)
    tb(slide, M, 5.18, W, 0.22, "SAME CONTRACT, DIFFERENT LOOK-AHEAD", font=MONO, size=8.8, bold=True, color=MUTED, align=PP_ALIGN.CENTER)
    variants = [
        ("STATIC", "No solve; capacity-proportional reference", MUTED),
        ("REACTIVE", "Responds after overload is observed", RED),
        ("LP", "Optimises the next decision only", BLUE),
        ("PRE-DRAIN", "LP activated by advance loss notice", AMBER),
        ("MPC", "LP across a horizon; apply first step, re-solve", PURPLE),
        ("ORACLE", "Whole future known; ceiling only", GREEN),
    ]
    vgap, vw = 0.12, 1.88
    for index, (name, body, accent) in enumerate(variants):
        left = M + index * (vw + vgap)
        rect(slide, left, 5.54, vw, 0.98, CARD, BORDER)
        rect(slide, left, 5.54, 0.04, 0.98, accent)
        tb(slide, left + 0.13, 5.68, vw - 0.26, 0.20, name, font=MONO, size=8.0, bold=True, color=accent, align=PP_ALIGN.CENTER)
        tb(slide, left + 0.13, 5.94, vw - 0.26, 0.39, body, size=7.9, color=BODY, align=PP_ALIGN.CENTER)


def deployment_architecture_slide(slide) -> None:
    chrome(
        slide,
        "REAL-WORLD DEPLOYMENT · OPERATING MODEL",
        "Deploy the evidence loop first; earn the right to automate",
        "03 · THE OPTIMIZER",
        "The production service runs continuously, but a recommendation changes the network only after an independent gate—and initially, an operator—accepts it.",
        source="ADR-011/012/019/024 · C-DOT first-drop advisory",
    )
    steps = [
        ("01", "OBSERVE", "UPF/SMF counters\nhealth · topology\nmaintenance notices", TEAL),
        ("02", "CANONICALISE", "map labels → IDs\nclose 10-min bucket\nflag gaps/resets", BLUE),
        ("03", "FORECAST", "p50/p90/p95\nnew-session demand\n10–80 min horizons", PURPLE),
        ("04", "OPTIMISE", "weights by scope\nStatic comparison\nconstraints + churn", AMBER),
        ("05", "CERTIFY", "freshness · capacity\neligibility · latency\nsolver + policy hash", GREEN),
        ("06", "ADVISE / ACT", "operator console\nversioned API + TTL\nverify or rollback", RED),
    ]
    sw, gap = 1.82, 0.18
    for index, (num, head, body, accent) in enumerate(steps):
        left = M + index * (sw + gap)
        rect(slide, left, 2.25, sw, 1.68, CARD, BORDER, radius=True)
        rect(slide, left, 2.25, sw, 0.08, accent)
        tb(slide, left + 0.12, 2.48, 0.30, 0.28, num, font=SERIF, size=16,
           color=accent, align=PP_ALIGN.CENTER)
        tb(slide, left + 0.43, 2.49, sw - 0.55, 0.20, head, font=MONO, size=7.8,
           bold=True, color=accent)
        tb(slide, left + 0.14, 2.91, sw - 0.28, 0.72, body, size=8.7,
           color=BODY, align=PP_ALIGN.CENTER, spacing=1.08)
        if index < len(steps) - 1:
            arrow(slide, left + sw + 0.025, 2.93, 0.13, 0.25, TEAL)

    # Closed-loop feedback strip.
    rect(slide, 1.03, 4.12, 11.30, 0.47, SOFT_BLUE, BORDER, radius=True)
    tb(slide, 1.20, 4.21, 10.96, 0.25,
       "NEXT CLOSED BUCKET → compare predicted vs realised load, requested vs realised shares, benefit vs Static, and any fallback; update coverage/drift monitors—not the protected release result.",
       size=9.2, bold=True, color=BLUE, align=PP_ALIGN.CENTER)

    stages = [
        (M, "1 · SHADOW", "No network write", "Generate signed JSON/CSV advisories, replay against recorded telemetry and score the counterfactual against Static.", TEAL),
        (4.47, "2 · SUPERVISED CANARY", "Human approval + bounded scope", "Apply to one agreed TAC/DNN or maintenance window, cap the weight change, set a TTL, verify the read-back and keep one-click rollback.", AMBER),
        (8.23, "3 · BOUNDED AUTOMATION", "Only after production gates", "Automatic apply requires confirmed IDs/semantics, calibrated envelopes, authenticated compare-and-swap, proven rollback, coverage/drift alarms and an audit trail.", GREEN),
    ]
    widths = [3.51, 3.51, 4.38]
    for index, (left, label, head, body, accent) in enumerate(stages):
        card(slide, left, 4.88, widths[index], 1.61, label, head, body,
             accent=accent, body_size=8.8, heading_size=10.8)
    rect(slide, M, 6.65, W, 0.18, SOFT_AMBER)
    tb(slide, M, 6.64, W, 0.18,
       "Production principle: advisory by default; actuation is a separately authorised capability.",
       font=MONO, size=8.1, bold=True, color=AMBER, align=PP_ALIGN.CENTER)


def risk_advisory_slide(slide) -> None:
    chrome(
        slide,
        "REAL-WORLD DEPLOYMENT · 10-MINUTE ADVISORY",
        "Publish weights with a risk envelope—not a naked recommendation",
        "03 · THE OPTIMIZER",
        "Risk is a vector of measurable quantities, not one opaque confidence score. The operator sees expected benefit, tail exposure, uncertainty, churn and data quality before choosing APPLY, HOLD or ROLLBACK.",
        source="advisory.json · ADR-009/011/024",
    )
    # Left: advisory envelope.
    rect(slide, M, 2.22, 4.00, 3.91, CARD, BORDER)
    rect(slide, M, 2.22, 0.06, 3.91, TEAL)
    tb(slide, 0.99, 2.40, 3.48, 0.22, "ADVISORY PACKET · EVERY 10 MINUTES", font=MONO,
       size=8.6, bold=True, color=TEAL)
    packet = [
        ("IDENTITY", "issued_at · target window · valid_until/TTL\ntelemetry watermark · model/policy hashes"),
        ("ACTION", "scope (TAC/DNN/slice) · current weights\nrecommended weights · maximum allowed change"),
        ("RISK VECTOR", "p50/p90/p95 load · p95 headroom\nbenefit vs Static · churn · quality flags"),
        ("DECISION", "RECOMMEND / HOLD / ROLLBACK\nreason codes · approvals · read-back status"),
    ]
    for index, (label, body) in enumerate(packet):
        top = 2.79 + index * 0.76
        tb(slide, 1.00, top, 0.91, 0.20, label, font=MONO, size=7.8, bold=True,
           color=TEAL)
        tb(slide, 1.93, top - 0.02, 2.49, 0.48, body, size=8.6, color=BODY,
           spacing=1.05)
        if index < len(packet) - 1:
            line(slide, 1.00, top + 0.58, 4.43, top + 0.58, RULE, 0.8)

    # Right: exact risk quantities.
    metrics = [
        ("TAIL HEADROOM", "min_u (C_u - L_u[p95]) / C_u", "How much capacity remains if demand reaches its p95 upper bound.", PURPLE),
        ("BENEFIT VS STATIC", "risk(static) − risk(advisory)", "Same-state reduction in overload severity; negative means the advisory is worse.", GREEN),
        ("FORECAST WIDTH", "(p95 - p50) / max(p50, eps)", "A dimensionless uncertainty band; wide forecasts demand more caution.", BLUE),
        ("POLICY CHURN", "0.5 Sum_g,u |w_new - w_current|", "How much routing changes. Limit both total churn and the largest group shift.", AMBER),
        ("DATA / MODEL HEALTH", "age · missing · resets · coverage · drift", "Publish components and reason codes; never hide them inside one score.", RED),
        ("REVERSIBILITY", "TTL · read-back · rollback verified", "An action is safer when it expires, can be verified, and can be reversed quickly.", TEAL),
    ]
    for index, (label, formula, body, accent) in enumerate(metrics):
        col, row = index % 2, index // 2
        left = 5.03 + col * 3.83
        top = 2.22 + row * 1.30
        rect(slide, left, top, 3.60, 1.12, CARD, BORDER)
        rect(slide, left, top, 0.045, 1.12, accent)
        tb(slide, left + 0.18, top + 0.13, 1.25, 0.18, label, font=MONO, size=7.4,
           bold=True, color=accent)
        tb(slide, left + 1.42, top + 0.10, 1.98, 0.24, formula, font=SERIF,
           size=10.4, color=INK, align=PP_ALIGN.RIGHT)
        tb(slide, left + 0.18, top + 0.45, 3.20, 0.50, body, size=8.2,
           color=BODY, spacing=1.05)

    # Operator decision bands.
    bands = [
        (5.03, 6.22, 2.28, "GREEN · RECOMMEND", "fresh + feasible + positive p95 headroom", GREEN),
        (7.43, 6.22, 2.28, "AMBER · APPROVE", "narrow margin or high churn; canary only", AMBER),
        (9.83, 6.22, 2.78, "RED · HOLD / ROLLBACK", "stale, invalid, unconfirmed or worse than Static", RED),
    ]
    for left, top, width, label, body, accent in bands:
        rect(slide, left, top, width, 0.48, CARD, BORDER)
        rect(slide, left, top, 0.05, 0.48, accent)
        tb(slide, left + 0.13, top + 0.08, width - 0.26, 0.15, label, font=MONO,
           size=6.8, bold=True, color=accent, align=PP_ALIGN.CENTER)
        tb(slide, left + 0.13, top + 0.25, width - 0.26, 0.14, body, size=6.9,
           color=BODY, align=PP_ALIGN.CENTER)
    rect(slide, M, 6.31, 4.00, 0.39, SOFT_RED, BORDER)
    tb(slide, 0.92, 6.38, 3.60, 0.23,
       "CURRENT FIRST DROP → RED / HOLD: mapping and SMF semantics unconfirmed; capacities uncalibrated; new-session arrivals absent.",
       size=7.7, bold=True, color=RED, align=PP_ALIGN.CENTER)


def add_notes(slide, text: str) -> None:
    slide.notes_slide.notes_text_frame.text = text


def renumber(prs) -> None:
    for index, slide in enumerate(prs.slides, 1):
        candidates = []
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            if shape.left > Inches(11.5) and shape.top > Inches(6.9):
                candidates.append(shape)
        for shape in candidates:
            set_existing_text(shape, f"{index:02d}", size=8.5, color=FAINT)
            shape.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT


def audit(prs) -> None:
    required = [
        "Every equation has one job",
        "Different simulators answer different 5G questions",
        "A 16-week extreme shard is a ~9.5-hour job",
        "Five model families: what each one is trying to learn",
        "How error, safety coverage and release gates are calculated",
        "Past telemetry goes in; demand distributions come out",
        "The optimizer turns demand and capacity into routing percentages",
        "Deploy the evidence loop first; earn the right to automate",
        "Publish weights with a risk envelope",
    ]
    all_text = "\n".join(slide_text(slide) for slide in prs.slides)
    missing = [item for item in required if item not in all_text]
    if missing:
        raise AssertionError(f"missing required content: {missing}")
    forbidden = [
        "LightGBM improved accuracy 11.6%",
        "LightGBM beats it by about 1.5 percentage points",
        "A version that knows the future removes 100% of overload.",
    ]
    present = [item for item in forbidden if item in all_text]
    if present:
        raise AssertionError(f"stale ambiguous wording remains: {present}")
    if prs.slide_width != Inches(13.333333) and abs(prs.slide_width - Inches(13.333333)) > 2000:
        raise AssertionError("unexpected presentation width")


def main() -> int:
    if not BACKUP.exists():
        shutil.copy2(DECK, BACKUP)
    # Rebuild from the preserved pre-feedback deck so rerunning the script is
    # deterministic and never compounds previous shape edits.
    prs = Presentation(BACKUP)

    rewrite_headline(prs)
    rewrite_area_slides(prs)

    add_after(prs, "How traffic is generated, mathematically", "Every equation has one job", equation_explainer)
    add_after(prs, "The realism layer added on top", "Different simulators answer different 5G questions", comparison_slide)
    add_after(prs, "How much headroom was left", "A 16-week extreme shard is a ~9.5-hour job", runtime_slide)
    add_after(prs, "Predicting how much traffic is coming", "Five model families: what each one is trying to learn", model_primer)

    rewrite_forecaster_flow(prs)
    add_after(prs, "Accuracy on 16 weeks of history", "How error, safety coverage and release gates are calculated", scoring_slide)
    rewrite_forecast_comparison(prs)
    rewrite_forecast_next_steps(prs)
    rewrite_optimizer_flow(prs)
    add_after(prs, "What we recommend running, and why", "Deploy the evidence loop first", deployment_architecture_slide)
    add_after(prs, "Deploy the evidence loop first", "Publish weights with a risk envelope", risk_advisory_slide)

    for title, note in {
        "Every equation has one job": "SAY: The generator is a pipeline. Scale and calendar determine expected arrivals; Poisson turns expectation into a count; lifetime creates persistence; summation creates UPF load; Little's Law and traffic conservation validate the output.",
        "Different simulators answer different 5G questions": "SAY: This is not a protocol-conformance or radio simulator. It complements those tools by testing long-horizon steering decisions and their safety evidence. Do not claim replacement of dsTest, PacketRusher, UERANSIM, 5G-LENA or Simu5G.\n\nPUBLIC SOURCES CHECKED 27 AUG 2026:\nhttps://www.mobileum.com/products/testing-assurance-observability/performance-and-functional-assurance/performance-testing\nhttps://github.com/HewlettPackard/PacketRusher\nhttps://github.com/aligungr/UERANSIM\nhttps://apps.nsnam.org/app/nr/\nhttps://simu5g.org/users-guide/overview",
        "A 16-week extreme shard is a ~9.5-hour job": "SAY: 9:31:46 is projected from a complete one-day measurement; reserve 12 hours. A single shard is not distributed across nodes. The 85.4-minute result is a measured 384-shard one-day campaign, not a 16-week timing.",
        "Five model families: what each one is trying to learn": "SAY: Ridge is a stable straight-line mapping. Boosted trees learn nonlinear thresholds. The regime ensemble switches models after a state becomes causally observable. Complexity must beat the moving-average baseline on every gate.",
        "How error, safety coverage and release gates are calculated": "SAY: 11.64% is a relative reduction in error. It is not forecast accuracy. The 15% threshold was frozen before selection, and LightGBM also failed the worst-slice guardrail, so retaining ridge follows the experiment contract.",
        "Deploy the evidence loop first": "SAY: The service may run every ten minutes from day one, but production write authority is a separate capability. Start with signed advisories and counterfactual scoring, then a supervised canary with TTL/read-back/rollback, and automate only after C-DOT confirms identifiers, policy semantics and capacity envelopes.",
        "Publish weights with a risk envelope": "SAY: Do not compress risk into one confidence percentage. Publish tail headroom, same-state benefit over Static, forecast width, policy churn, data/model health and reversibility. The existing first C-DOT drop is a useful RED/HOLD example because the identity mapping, capacity values, session arrivals and SMF semantics are not yet production-safe.",
    }.items():
        add_notes(find_slide(prs, title), note)

    renumber(prs)
    audit(prs)
    prs.save(DECK)
    print(f"updated={DECK}")
    print(f"backup={BACKUP}")
    print(f"slides={len(prs.slides)}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
