"""Make the v4 C-DOT PowerPoint internally consistent for the office review.

The earlier experiment slides remain unchanged as an audit trail. This script only
updates summary, section-introduction, decision, inventory, and conclusion language
that still presented the pre-v4 conclusion as current.
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation


DECK = Path(__file__).resolve().parent / "CDOT_UPF_Steering_Evidence_Review.pptx"


def replace_text(slide, old: str, new: str) -> bool:
    # Prefer a whole shape or whole paragraph. Short labels such as "NO" must
    # never be treated as substrings of words such as "NODES".
    for shape in slide.shapes:
        if shape.has_text_frame and shape.text == old:
            lines = new.split("\n")
            paragraphs = shape.text_frame.paragraphs
            while len(paragraphs) < len(lines):
                shape.text_frame.add_paragraph()
                paragraphs = shape.text_frame.paragraphs
            for index, paragraph in enumerate(paragraphs):
                line = lines[index] if index < len(lines) else ""
                if paragraph.runs:
                    paragraph.runs[0].text = line
                    for extra in paragraph.runs[1:]:
                        extra.text = ""
                elif line:
                    paragraph.add_run().text = line
            return True
    if len(old) <= 3:
        return False
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        if old in shape.text and "\n" in old:
            replacement = shape.text.replace(old, new)
            lines = replacement.split("\n")
            paragraphs = shape.text_frame.paragraphs
            while len(paragraphs) < len(lines):
                shape.text_frame.add_paragraph()
                paragraphs = shape.text_frame.paragraphs
            for index, paragraph in enumerate(paragraphs):
                line = lines[index] if index < len(lines) else ""
                if paragraph.runs:
                    paragraph.runs[0].text = line
                    for extra in paragraph.runs[1:]:
                        extra.text = ""
                elif line:
                    paragraph.add_run().text = line
            return True
        for paragraph in shape.text_frame.paragraphs:
            joined = "".join(run.text for run in paragraph.runs)
            if old in joined and paragraph.runs:
                paragraph.runs[0].text = joined.replace(old, new)
                for extra in paragraph.runs[1:]:
                    extra.text = ""
                return True
    return False


# Ordered PowerPoint slide number, old text, replacement text.
REPLACEMENTS = [
    (6, "RUNS ON 12\nCLUSTER PASSDES", "RUNS ON 12\nCLUSTER NODES"),
    (6, "NO", "PASS"),
    (6, "MPC CONTROLLER\nNOT RELEASED", "MPC · DECLARED EVENTS\nSHADOW CANDIDATE"),
    (6, "NO", "PASS"),
    (6, "PRE-DRAIN\nNOT RELEASED", "PRE-DRAIN · DECLARED\nEVENTS ONLY"),
    (6, "SIMPLE", "STATIC"),
    (6, "THE RULE WE\nSTILL USE", "DEFAULT OUTSIDE\nDECLARED EVENTS"),
    (6, "WHAT DID NOT", "WHAT THE AUDIT CHANGED"),
    (6, "Worst-case safety", "Mixed-event exposure remains"),
    (6, "Best candidate failed the worst-day limit by 0.37 points.",
        "Pre-drain can regress when declared loss is followed by a surprise; MPC cuts that exposure."),
    (6, "Decision speed", "Faster is not automatically better"),
    (6, "225–970 ms per decision; only one candidate met the 500 ms budget.",
        "Two-minute cadence bought 0.02 points for 3.7× the routing churn."),
    (42, "Seven controller designs, 28 configurations, 588 head-to-head tests. The optimizer works: it solves, it validates, it acts. It just does not reliably beat a simple fixed rule. The reason turns out to be about the lever we have, not the algorithm.",
        "Experiments 1–6 show how the first designs failed. Experiment 7 records the evaluator audit and the corrected result: predictive steering wins for declared capacity loss, ties static on pure surprises, and remains shadow-only."),
    (44, "IN PRODUCTION", "REFERENCE DEFAULT"),
    (54, "NOBODY FILLED A ROW", "NO EARLY CANDIDATE FILLED A ROW"),
    (63, "The simple rule stays the default. Guarded pre-drain or cohort MPC is enabled for declared maintenance windows only, and falls back to the simple rule everywhere else.",
        "Static is the reference default. Evaluate guarded pre-drain or cohort MPC in shadow mode for declared maintenance windows only; no live policy publication is authorized."),
    (63, "It wins on the metric we chose in advance", "Static remains the safe default"),
    (63, "Across 128 production test days it is best on uplink overload, downlink overload, uplink drops and downlink drops — all four at once.",
        "Outside a declared maintenance window, publish the capacity-proportional static policy. Pure surprises are verified bit-exact ties."),
    (63, "Spreading is structurally safer", "Declared events unlock predictive value"),
    (63, "When each group has three to six eligible boxes, sessions last hours, and you cannot move them afterwards, spreading beats concentrating.",
        "On held-out declared-event pairs, predictive steering removes overload without a single regression across 473 informative cases."),
    (63, "It has no bad days", "MPC is the conservative shadow candidate"),
    (63, "No worst-case loss, no validation to fail, no solver to time out, no data to go stale, no forecast to be wrong during a surge.",
        "Cohort MPC gives up some gain for roughly one tenth of pre-drain's collateral harm and routing churn."),
    (63, "The clever controllers stay available, but supervised", "Fail closed on uncertainty"),
    (63, "MPC and pre-drain remain implemented and can run in shadow mode. They give up control the moment telemetry looks uncertain, their data is stale, or a capacity limit would be exceeded.",
        "Any stale telemetry, invalid plan, solver failure, undeclared event, or uncertain envelope retains the last safe static policy."),
    (63, "100", "STATIC"),
    (63, "SIMPLE RULE\nBASELINE", "DEFAULT\nPOLICY"),
    (63, "105–117", "+24.0%"),
    (63, "COHORT\nMPC", "HELD-OUT\nGAIN"),
    (63, "198–221", "0 / 473"),
    (63, "REACTIVE\nCONTROLLER", "DECLARED-EVENT\nREGRESSIONS"),
    (66, "588", "25 000"),
    (66, "CONTROLLER\nCOMPARISONS", "PAIRED SIMULATED\nDAYS"),
    (66, "28", "160"),
    (66, "CONFIGURATIONS\nTESTED", "V4 CONFIGURATIONS\nSCREENED"),
    (66, "125", "36"),
    (66, "BLIND SURVIVAL\nTRIALS", "PASSED EVERY\nGATE"),
    (66, "174", "3 / 4"),
    (66, "CODE TESTS\nPASSING", "SURVIVED FRESH-SEED\nVALIDATION"),
    (67, "Pure surprises remain bit-exact ties.e stays.", "Pure surprises remain bit-exact ties."),
    (68, "Yes turns 1.6% movable traffic into a solvable problem", "Yes turns about 3% movable traffic into a more tractable problem"),
]


def main() -> int:
    presentation = Presentation(DECK)
    missed: list[tuple[int, str]] = []
    for slide_number, old, new in REPLACEMENTS:
        slide = presentation.slides[slide_number - 1]
        if not replace_text(slide, old, new) and not any(
            shape.has_text_frame and new in shape.text for shape in slide.shapes
        ):
            missed.append((slide_number, old))
    if missed:
        for slide_number, old in missed:
            print(f"MISSED slide {slide_number}: {old!r}")
        return 1
    presentation.save(DECK)
    print(f"Updated {len(REPLACEMENTS)} v4 consistency statements in {DECK.name}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
